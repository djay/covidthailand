import concurrent.futures
import datetime
import json
import os
import pickle
import random
import re
import sys
import tempfile
import urllib.parse
from io import StringIO
from itertools import compress
from itertools import cycle
from pathlib import Path

import camelot
import dateutil
import pandas as pd
import pythainlp
import requests
import urllib3
from bs4 import BeautifulSoup
from loguru import logger
from pptx import Presentation
try:
    from pytwitterscraper import TwitterScraper
except:
    TwitterScraper = None
from requests.adapters import HTTPAdapter
from requests.adapters import Retry
from requests.exceptions import RequestException
from requests.exceptions import Timeout
from tika import config
from tika import parser
from webdav3.client import Client
from xlsx2csv import Xlsx2csv
#from proxyscrape import create_collector


# https://proxybroker.readthedocs.io/en/latest/

# collector = create_collector('my-collector', 'http')
# collectors = create_collector('my-collectors', 'https')

# import asyncio
# from proxybroker import Broker


# proxy = None
# async def show(proxies):
#     global proxy
#     proxy = await proxies.get()

# proxies = asyncio.Queue()
# broker = Broker(proxies)
# tasks = asyncio.gather(
#     broker.find(types=['HTTP', 'HTTPS'], limit=1),
#     show(proxies))

# loop = asyncio.get_event_loop()
# loop.run_until_complete(tasks)


urllib3.disable_warnings()

CHECK_NEWER = bool(os.environ.get("CHECK_NEWER", False))
USE_CACHE_DATA = os.environ.get('USE_CACHE_DATA', False) == 'True'
MAX_DAYS = int(os.environ.get("MAX_DAYS", 1 if USE_CACHE_DATA else 0))

NUM_RE = re.compile(r"\d+(?:\,\d+)*(?:\.\d+)?")
INT_RE = re.compile(r"\d+(?:\,\d+)*")
NUM_OR_DASH = re.compile(r"([0-9\,\.]+|-)-?")

requests.adapters.DEFAULT_RETRIES = 3  # for other tools that use requests internally
RETRY = Retry(
    total=3, backoff_factor=1
)  # should make it more reliable as ddc.moph.go.th often fails


DEFAULT_TIMEOUT = 5  # seconds


class TimeoutHTTPAdapter(HTTPAdapter):
    def __init__(self, *args, **kwargs):
        self.timeout = DEFAULT_TIMEOUT
        if "timeout" in kwargs:
            self.timeout = kwargs["timeout"]
            del kwargs["timeout"]
        super().__init__(*args, **kwargs)

    def send(self, request, **kwargs):
        timeout = kwargs.get("timeout")
        if timeout is None:
            kwargs["timeout"] = self.timeout
        return super().send(request, **kwargs)


def fix_timeouts(s, timeout=None):
    if timeout is not None:
        adapter = TimeoutHTTPAdapter(max_retries=RETRY, timeout=timeout)
    else:
        adapter = TimeoutHTTPAdapter(max_retries=RETRY)
    s.mount("http://", adapter)
    s.mount("https://", adapter)


s = requests.Session()
fix_timeouts(s)


# do any tika install now before we start the run and use multiple processes
config.getParsers()


def today() -> datetime.datetime:
    """Return today's date and time"""
    return datetime.datetime.today()


def formatter(log_entry):
    """
    Formats a log entry allowing for multiple entries per line
    Required for progress indicators...
    """
    end = log_entry["extra"].get("end", "\n")
    return "[{time}] {message}" + end + "{exception}"


# first we clear the multi-process-naive default logger
logger.remove()
# then we add an MP-aware one instead
logger.add(sys.stderr, format=formatter, enqueue=True)


####################
# Extraction helpers
#####################
def parse_file(filename, html=False, paged=True, remove_corrupt=True):
    pages_txt = []

    # Read PDF file
    data = parser.from_file(filename, xmlContent=True)
    if not data or not data["content"] and remove_corrupt:
        # file is corrupt. Delete is so can get redownloaded
        os.remove(filename)
        return "" if not paged else []
    xhtml_data = BeautifulSoup(data["content"], features="lxml")
    if html and not paged:
        return xhtml_data
    pages = xhtml_data.find_all("div", attrs={"class": ["page", "slide-content"]})
    if not pages:
        if not paged:
            return repr(xhtml_data)
        else:
            return [repr(xhtml_data)]

    # TODO: slides are divided by slide-content and slide-master-content rather than being contained
    for i, content in enumerate(pages):
        # Parse PDF data using TIKA (xml/html)
        # It's faster and safer to create a new buffer than truncating it
        # https://stackoverflow.com/questions/4330812/how-do-i-clear-a-stringio-object
        _buffer = StringIO()
        _buffer.write(str(content))
        parsed_content = parser.from_buffer(_buffer.getvalue())
        if parsed_content["content"] is None:
            pages_txt.append("")
        else:
            # Add pages
            text = parsed_content["content"].strip()
            if html:
                pages_txt.append(repr(content))
            else:
                pages_txt.append(text)
    if paged:
        return pages_txt
    else:
        return '\n\n\n'.join(pages_txt)


def get_next_numbers(content, *matches, debug=False, before=False, remove=0, ints=True, until=None, return_rest=True, return_until=False, require_until=False, dash_as_zero=False, thainorm=False, asserted=False):
    """
    returns the numbers that appear immediately before or after the string(s) in 'matches',
    optionally up through 'until', that are found in the parsed PDF string 'content'
    """
    if thainorm:
        content = pythainlp.util.remove_tonemark(pythainlp.util.normalize(content))
        until = pythainlp.util.remove_tonemark(pythainlp.util.normalize(until))
        matches = [pythainlp.util.remove_tonemark(pythainlp.util.normalize(match)) for match in matches]

    if len(matches) == 0:
        matches = [""]
    for match in matches:
        if type(match) == str:
            match = re.compile(f"({match})")
        ahead, *behind = match.split(content, 1) if match else ("", "", content)
        if not behind:
            continue
        matched, *behind = behind
        behind = "".join(behind)
        found = ahead if before else behind
        if until is not None:
            found, *rest = re.split(until, found, 1)  # TODO: how to put it back together if behind=True?
            if not rest and require_until:
                # in this case return nothing since end didn't find a match
                continue
            rest = until + (rest[0] if rest and rest[0] else "")
        else:
            rest = ""
        if dash_as_zero:
            found = found.replace(r'-', '0')
        numbers = (INT_RE if ints else NUM_RE).findall(found)
        numbers = [n.replace(",", "") for n in numbers]
        numbers = [int(n) if ints else float(n) for n in numbers if n]
        numbers = numbers if not before else list(reversed(numbers))
        if remove:
            behind = (INT_RE if ints else NUM_RE).sub("", found, remove)
        if return_until:
            return numbers, found
        elif return_rest:
            return numbers, matched + " " + rest + behind
        else:
            return numbers
    if debug and matches:
        logger.info("Couldn't find '{}'", match)
        logger.info(content)
    if asserted:
        assert False, f"None of {matches} in: {content}"
    if return_rest or return_until:
        return [], content
    else:
        return []


def get_next_number(content, *matches, default=None, remove=False, before=False, until=None, return_rest=True,
                    require_until=False, dash_as_zero=False, thainorm=False, asserted=False):
    num, rest = get_next_numbers(content, *matches, remove=1 if remove else 0, before=before, until=until,
                                 require_until=require_until, dash_as_zero=dash_as_zero, thainorm=thainorm, asserted=asserted)
    num = num[0] if num else default
    if return_rest:
        return num, rest
    else:
        return num


def toint(s):
    return int(s.replace(',', '')) if s else None


def slide2text(slide):
    text = ""
    if slide.shapes.title:
        text += slide.shapes.title.text
    for shape in slide.shapes:
        if shape.has_text_frame:
            # for p in shape.text_frame:
            text += "\n" + shape.text
    return text


def pptx2chartdata(file):

    def find_charts(shape, i):
        if not shape.has_chart:
            for s in getattr(shape, 'shapes', []):  # Group shapes
                yield from find_charts(s, i)
            return
        chart = shape.chart
        if chart is None:
            return
        title = chart.chart_title.text_frame.text if chart.has_title else ""
        series = dict([(s.name, s.values) for s in chart.series])
        yield chart, title, series, i

    prs = Presentation(file)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            yield from find_charts(shape, i)


####################
# Download helpers
####################
def resume_from(file, remote_date, check=True, size=0, appending=False):
    if type(remote_date) == str:
        remote_date = dateutil.parser.parse(remote_date)

    if not os.path.exists(file):
        logger.info("Missing: {}", file)
        return 0
    else:
        fdate = datetime.datetime.fromtimestamp(os.path.getmtime(file)).astimezone()
        resume_pos = os.stat(file).st_size

    if resume_pos == 0:
        # probably something went wrong. redownload
        return 0
    elif size and size != resume_pos:
        return 0
    elif not check:
        return -1
    elif remote_date is None:
        return 0  # TODO: should we always keep cached?
    elif size and size == resume_pos and remote_date <= fdate:
        # it's the same, don't redownload
        logger.info("Cached Unmodified: {} <= {} {}b: {}", remote_date, fdate, size, file)
        return -1
    elif appending and size:
        if resume_pos < size:
            return resume_pos
        elif resume_pos > size:
            # redownload it
            return 0
        elif remote_date > fdate:
            # size is the same but says updated? redownload it to be sure
            return 0
        else:
            return -1
    elif remote_date > fdate:
        return 0
    elif size and resume_pos != size:
        return 0
    else:
        # same size and date so keep what we have
        return -1


def is_cutshort(file, modified, check):

    if type(modified) == str:
        modified = dateutil.parser.parse(modified)
    if not check and MAX_DAYS and modified and (datetime.datetime.today().astimezone()
                                                - modified).days > MAX_DAYS and os.path.exists(file):
        logger.info("Reached MAX_DAYS={}", MAX_DAYS)
        return True
    return False


def url2filename(url, strip_version=False):
    file = sanitize_filename(url.rsplit("/", 1)[-1])
    if strip_version and '.' in file:
        file = ".".join(file.split(".")[:2])
    return file


def links_html_namer(url, _):
    return "-".join(url.split("/")[2:]) + ".html"


def web_links(*index_urls, ext=".pdf", dir="inputs/html", match=None, filenamer=links_html_namer, check=True, timeout=None, proxy=False, ignore_errors=False):
    def is_ext(a):
        return len(a.get("href").rsplit(ext)) == 2 if ext else True

    def is_match(a):
        return a.get("href") and is_ext(a) and (match.search(a.get_text(strip=True)) if match else True)

    for file, index, index_url in web_files(*index_urls, dir=dir, check=check, filenamer=filenamer, timeout=timeout, proxy=proxy, threads=1):
        if ignore_errors and file is None:
            continue
        assert file is not None, f"Problem accessing {index_url}"
        soup = parse_file(file, html=True, paged=False)
        links = (urllib.parse.urljoin(index_url, a.get('href')) for a in soup.find_all('a') if is_match(a))
        for link in links:
            yield link


def web_files(*urls, dir=os.getcwd(), check=CHECK_NEWER, strip_version=False, appending=False, filenamer=url2filename, timeout=None, proxy=False, threads=5):
    """if check is None, then always download"""
    s = requests.Session()
    if timeout is None:
        timeout = 10
        # We only want retries under normal conditions
    fix_timeouts(s, timeout)

    def get_file(url, i, check):
        file = filenamer(url, strip_version)
        file = os.path.join(tempfile.gettempdir() if dir is None else dir, file)
        os.makedirs(os.path.dirname(file), exist_ok=True)
        resumable = False
        size = None
        #verify = "ddc.moph.go.th" not in url
        verify = True

        remove = False
        if check or MAX_DAYS:
            proxies = next(proxies_itor, None) if proxy else None
            try:
                r = s.head(url, timeout=timeout, verify=verify, proxies=proxies)
                modified = r.headers.get("Last-Modified")
                if r.headers.get("content-range"):
                    pre, size = r.headers.get("content-range").split("/")
                    size = int(size)
                    assert "bytes" in pre
                else:
                    size = int(r.headers.get("content-length", 0))
                resumable = r.headers.get('accept-ranges') == 'bytes' and check and size > 0
            except (Timeout, RequestException):
                modified = None
        else:
            modified = None
        if i > 0 and is_cutshort(file, modified, check):
            return None, None, url
        err = ""
        if (resume_byte_pos := resume_from(file, modified, check, size, appending)) < 0:
            if check:
                # logger.info("Unmodified: {}: using cache. {} {}", file, modified, size)
                pass
        else:
            # go back 10% in case end of data changed (e.g csv)
            resume_byte_pos = int(resume_byte_pos * 0.95) if resumable else 0
            resume_header = {'Range': f'bytes={resume_byte_pos}-'} if resumable else {}

            proxies = next(proxies_itor, None) if proxy else None
            try:
                # handle resuming based on range requests - https://stackoverflow.com/questions/22894211/how-to-resume-file-download-in-python
                # Speed up covid-19 download a lot, but might have to jump back to make sure we don't miss data.
                r = s.get(url, timeout=timeout, stream=True, headers=resume_header, allow_redirects=True,
                          verify=verify, proxies=proxies)
            except (Timeout, RequestException) as e:
                err = str(e)
                r = None
            #if type(check) == int and check > 0:
            #    check -= 1  # HACK: using int as Boolean above
            if r is None or r.status_code >= 300:
                err = f"bad response {r.status_code}, {r.content}" if r is not None else err
                if not os.path.exists(file):
                    logger.info("Error downloading: {}: skipping. {}", file, err)
                    return None, None, url
                logger.info("Error downloading: {}: using cache. {}", file, err)
            else:
                logger.bind(end="").opt(raw=True).info("Download: {} {}", file, modified)
                os.makedirs(os.path.dirname(file), exist_ok=True)
                mode = "w+b" if resume_byte_pos > 0 else "wb"
                with open(file, mode) as f:
                    f.seek(resume_byte_pos, 0)
                    # TODO: handle timeouts happening below since now switched to streaming
                    try:
                        for chunk in r.iter_content(chunk_size=2 * 1024 * 1024):
                            if chunk:  # filter out keep-alive new chunks
                                f.write(chunk)
                                logger.bind(end="").opt(raw=True).info(".")
                    except (Timeout, RequestException) as e:
                        if resumable:
                            # TODO: should we revert to last version instead?
                            logger.opt(raw=True).info("Error downloading: {}: resumable file incomplete {}", file, str(e))
                            return None, None, url
                        else:
                            logger.opt(raw=True).info("Error downloading: {}: skipping. {}", file, str(e))
                            remove = True
                logger.opt(raw=True).info("\n")
            logger.bind(end="\n")
        if remove:
            os.remove(file)  # if we leave it without check it will never get fixed
            return None, None, url
        with open(file, "rb") as f:
            content = f.read()
        if dir is None:
            os.remove(file)
        # i += 1
        return file, content, url

    with concurrent.futures.ThreadPoolExecutor(max_workers=threads) as executor:
        if type(check) == int:
            checks = [True] * check + [False] * (len(urls) - check)
        else:
            checks = [check] * len(urls)
        yield from executor.map(get_file, urls, range(len(urls)), checks)


def sanitize_filename(filename):
    return filename.translate(str.maketrans({"*": "_", "?": "_", ":": "_", "\\": "_", "<": "_", ">": "_", "|": "_"}))
    # Windows Filename Compatibility: '?*:<>|'


def local_files(ext=".pdf", dir=os.getcwd()):
    client_list = [
        (
            f"{file}",
            datetime.datetime.fromtimestamp(file.stat().st_ctime, tz=datetime.timezone.utc)
        )
        for file in Path(dir).glob(f'**/*{ext}') if file.is_file()
    ]
    # important we get them sorted newest files first as we only fill in NaN from each additional file
    files = sorted(
        client_list,
        key=lambda info: info[1],
        reverse=True,
    )
    i = 0
    for info in files:
        file = info[0].split("/")[-1]
        if not any([ext == file[-len(ext):] for ext in ext.split()]):
            continue
        target = os.path.join(dir, file)
        os.makedirs(os.path.dirname(target), exist_ok=True)
        if i > 0 and is_cutshort(target, info[1], False):
            break

        def do_dl(target=target): return target
        i += 1
        yield target, do_dl


def dav_files(url, username=None, password=None,
              ext=".pdf .pptx", dir=os.getcwd()):

    options = {
        "webdav_hostname": url,
        "webdav_login": username,
        "webdav_password": password,
    }
    client = Client(options)
    fix_timeouts(client.session)
    use_cache = False
    try:
        client_list = [(info["path"], info["modified"]) for info in client.list(get_info=True)]
    except (Exception):
        client_list = [
            (
                f"{file}",
                f"{datetime.datetime.fromtimestamp(file.stat().st_mtime, tz=datetime.timezone.utc):%a, %d %b %Y %H:%M:%S %Z}"
            )
            for file in Path(dir).glob('**/*') if file.is_file()
        ]
        use_cache = True
    # important we get them sorted newest files first as we only fill in NaN from each additional file
    files = sorted(
        client_list,
        key=lambda info: info[1],
        reverse=True,
    )
    i = 0
    for info in files:
        file = info[0].split("/")[-1]
        if not any([ext == file[-len(ext):] for ext in ext.split()]):
            continue
        target = os.path.join(dir, file)
        os.makedirs(os.path.dirname(target), exist_ok=True)
        if i > 0 and is_cutshort(target, info[1], False):
            break
        if not use_cache and resume_from(target, info[1]) >= 0:
            def do_dl(file=file, target=target):
                client.download_file(file, target)
                return target
        else:
            def do_dl(target=target): return target
        i += 1
        yield target, do_dl


#################
# Twitter helpers
#################
def parse_tweet(tw, tweet, found, *matches):
    """if tweet contains any of matches return its text joined with comments by the same person
    that also match (and contain [1/2] etc)"""
    if not any_in(tweet.get('text', tweet.get("comment", "")), *matches):
        return ""
    text = tw.get_tweetinfo(tweet['id']).contents['text']
    if any(text in t for t in found):
        return ""
    # TODO: ensure tweets are [1/2] etc not just "[" and by same person
    if "[" not in text:
        return text
    for t in sorted(tw.get_tweetcomments(tweet['id']).contents, key=lambda t: t['id']):
        rest = parse_tweet(tw, t, found + [text], *matches)
        if rest and rest not in text:
            text += " " + rest
    return text


def get_tweets_from(userid, datefrom, dateto, *matches):
    """return tweets from single person that match, merging in followups of the form [1/2]. Caches to speed up"""

    filename = os.path.join("inputs", "tweets", f"tweets2_{userid}.pickle")
    os.makedirs("inputs/tweets", exist_ok=True)
    try:
        with open(filename, "rb") as fp:
            tweets = pickle.load(fp)
    except (IOError, EOFError, OSError, pickle.PickleError, pickle.UnpicklingError) as e:
        logger.info('Error detected when attempting to load the pickle file: {}, setting an empty \'tweets\' dictionary', e)
        tweets = {}
    for date, tweet_list in tweets.items():
        fixed = []
        for tweet in tweet_list:
            text, url = (tweet, None) if type(tweet) == str else tweet
            fixed.append((text, (url if url else None)))
        tweets[date] = fixed
    latest = max(tweets.keys()) if tweets else None
    return tweets  # seems to be blocking and not useful new data anymore 2021-01-22
    tw = TwitterScraper()
    if latest and dateto and latest >= (datetime.datetime.today() if not dateto else dateto).date():
        return tweets
    for limit in [50, 2000, 20000]:
        logger.info("Getting {} tweets", limit)
        try:
            resp = tw.get_tweets(userid, count=limit).contents
        except Exception:
            # Either requests exception or intermittent Exception("ID User Not Found!")
            resp = []
        for tweet in sorted(resp, key=lambda t: t['id']):
            date = tweet['created_at'].date()
            url = tweet['urls'][0]['url'] if tweet['urls'] else f"https://twitter.com/{userid}/status/{tweet['id']}"
            text = parse_tweet(tw, tweet, tweets.get(date, []), *matches)
            if text:
                tweets[date] = tweets.get(date, []) + [(text, url)]

        earliest = min(tweets.keys())
        latest = max(tweets.keys())
        logger.info("got tweets {} to {} {}", earliest, latest, len(tweets))
        if earliest <= datefrom.date():  # TODO: ensure we have every tweet in sequence?
            break
        else:
            logger.info("Retrying: Earliest {}", earliest)
    with open(filename, "wb") as fp:
        pickle.dump(tweets, fp)
    return tweets


#################
# String helpers
#################
def remove_prefix(text: str, *prefixes: str) -> str:
    """Removes the prefix of a string"""
    for prefix in prefixes:
        if text.startswith(prefix):
            text = text[len(prefix):]
    return text


def remove_suffix(text: str, *suffixes: str) -> str:
    """Removes the suffix of a string"""
    for suffix in suffixes:
        if suffix and text.endswith(suffix):
            text = text[:-len(suffix)]
    return text


def seperate(seq, condition):
    a, b = [], []
    for item in seq:
        (a if condition(item) else b).append(item)
    return a, b


def split(seq, condition, maxsplit=0):
    """Similar to str.split except works on lists of lines. e.g. split([1,2,3,4], lambda x: x==2) -> [[1],[2],[3,4]]"""
    run = []
    last = False
    splits = 0
    for i in seq:
        if (maxsplit and splits >= maxsplit) or bool(condition(i)) == last:
            run.append(i)
        else:
            splits += 1
            yield run
            run = [i]
            last = not last
    yield run


def pairwise(lst):
    """Takes a list and turns them into pairs of tuples, e.g. [1,2,3,4] -> [[1,2],[3,4]]"""
    lst = list(lst)
    return list(zip(compress(lst, cycle([1, 0])), compress(lst, cycle([0, 1]))))


# def nwise(iterable, n=2):
#     iters = tee(iterable, n)
#     for i, it in enumerate(iters):
#         next(islice(it, i, i), None)
#     return zip(*iters)

def parse_numbers(lst):
    return [float(i.replace(",", "")) if i != "-" else 0 for i in lst]


def any_in(target, *matches):
    return any((str(m) in target) if type(m) != re.Pattern else m.search(target) for m in matches)


def all_in(target, *matches):
    return all((m in target) if type(m) != re.Pattern else m.search(target) for m in matches)


def strip(lst):
    lst = [i.strip() for i in lst]
    return [i for i in lst if i]


def unique_values(iterable):
    it = iter(iterable)
    seen = set()
    for item in it:
        if item in seen:
            continue
        seen.add(item)
        yield item


def replace_matcher(matches, replacements=None):
    if replacements is None:
        replacements = matches

    def replace_match(item):
        for m, r in zip(matches, replacements):
            if re.search(m, item, re.IGNORECASE):
                return r
        return item
    return replace_match


def camelot_cache(file, page_num, process_background=False, table=0):
    fname = f"{os.path.basename(file)}.{page_num}.{table}.{process_background}.json"
    os.makedirs("inputs/camelot", exist_ok=True)
    cache_file = os.path.join("inputs/camelot", fname)
    if os.path.exists(cache_file):
        return pd.read_json(cache_file)
    else:
        try:
            tables = camelot.read_pdf(file, pages=str(page_num), process_background=process_background)
        except ZeroDivisionError:
            tables = []
        if len(tables) < table + 1:
            with open(cache_file, "w") as fp:
                json.dump(None, fp)
            return None
        else:
            tables[table].df.to_json(cache_file)
            return tables[table].df


def read_excel(path: str, sheet_name: str = None) -> pd.DataFrame:
    buffer = StringIO()
    Xlsx2csv(path, outputencoding="utf-8", sheet_name=sheet_name).convert(buffer)
    buffer.seek(0)
    df = pd.read_csv(buffer)
    return df


def get_proxy():
    url = "https://raw.githubusercontent.com/mertguvencli/http-proxy-list/main/proxy-list/data-with-geolocation.json"
    url = "https://raw.githubusercontent.com/jetkai/proxy-list/main/online-proxies/json/proxies-advanced.json"
    try:
        data = requests.get(url, timeout=60).json()
    except requests.exceptions.RequestException:
        return
    random.shuffle(data)

    def to_proxies(d):
        for p in d['protocols']:
            _type = "http" if "http" in p['type'] else p['type']
            yield {
                "http": f"{_type}://{d['ip']}:{p['port']}",
                "https": f"{_type}://{d['ip']}:{p['port']}",
            }

        # if "http" in [p["type"] for p in d['protocols']]:
        #     return {
        #         # p["type"]: f"http{'s' if p['tls'] else ''}://{d['ip']}:{p['port']}" for p in d['protocols']
        #         p["type"]: f"http://{d['ip']}:{p['port']}" for p in d['protocols']
        #     }
        # else:
        #     p = d['protocols'][0]
        #     return {
        #         "http": f"{p['type']}://{d['ip']}:{p['port']}",
        #         "https": f"{p['type']}://{d['ip']}:{p['port']}",
        #     }
    proxies = [{
        "http": "socks4://127.0.0.1:8080",
        "https": "socks4://127.0.0.1:8080"
    },
        {
        "http": "socks4://127.0.0.1:9050",
        "https": "socks4://127.0.0.1:9050"
    }] + [p for d in data if d['location']['isocode'] == "TH" for p in to_proxies(d)]

    def test_proxy(proxies):
        try:
            if requests.head("https://ddc.moph.go.th", proxies=proxies, timeout=15).status_code < 400:
                return proxies
        except requests.exceptions.RequestException:
            # print("x", end="")
            pass
        # logger.info(f"Failed Test proxy {proxies}")

    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as executor:
        for future in concurrent.futures.as_completed(executor.submit(test_proxy, p) for p in proxies):
            proxies = future.result()
            if proxies is None:
                continue
            logger.info(f"Pass Test proxy {proxies}")
            yield proxies
            while test_proxy(proxies):
                logger.info(f"Pass ReTest proxy {proxies}")
                yield proxies


proxies_itor = get_proxy()
