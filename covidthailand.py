# coding=utf8
import datetime
import dateutil
import difflib
from io import StringIO
from itertools import islice, compress, cycle
import json
import os
import pathlib
import pickle
import re
from typing import Sequence, Union, List, Callable
import urllib.parse

from bs4 import BeautifulSoup
import camelot
import matplotlib
import matplotlib.pyplot as plt
from matplotlib.colors import ListedColormap, LinearSegmentedColormap
from matplotlib.ticker import FuncFormatter
from matplotlib.pyplot import cycler
import matplotlib.cm
import numpy as np
import pandas as pd
from pptx import Presentation
from pytwitterscraper import TwitterScraper
import requests
from requests.adapters import HTTPAdapter, Retry
from tika import parser
from webdav3.client import Client


CHECK_NEWER = bool(os.environ.get("CHECK_NEWER", False))

requests.adapters.DEFAULT_RETRIES = 5  # for other tools that use requests internally
s = requests.Session()
RETRY = Retry(
    total=10, backoff_factor=1
)  # should make it more reliable as ddc.moph.go.th often fails
s.mount("http://", HTTPAdapter(max_retries=RETRY))
s.mount("https://", HTTPAdapter(max_retries=RETRY))


#################
# String helpers
#################
def remove_prefix(text: str, prefix: str) -> str:
    """Removes the prefix of a string"""
    if text.startswith(prefix):
        return text[len(prefix):]
    else:
        return text


def remove_suffix(text: str, suffix: str) -> str:
    """Removes the suffix of a string"""
    if text.endswith(suffix):
        return text[:-len(suffix)]
    else:
        return text


###############
# Date helpers
###############
def today() -> datetime.datetime:
    """Return today's date and time"""
    return datetime.datetime.today()


THAI_ABBR_MONTHS = [
    "ม.ค.",
    "ก.พ.",
    "มี.ค.",
    "เม.ย.",
    "พ.ค.",
    "มิ.ย.",
    "ก.ค.",
    "ส.ค.",
    "ก.ย.",
    "ต.ค.",
    "พ.ย.",
    "ธ.ค.",
]
THAI_FULL_MONTHS = [
    "มกราคม",
    "กุมภาพันธ์",
    "มีนาคม",
    "เมษายน",
    "พฤษภาคม",
    "มิถุนายน",
    "กรกฎาคม",
    "สิงหาคม",
    "กันยายน",
    "ตุลาคม",
    "พฤศจิกายน",
    "ธันวาคม",
]


def file2date(file):
    "return date of either for '10-02-21' or '100264'"
    file = os.path.basename(file)
    file, *_ = file.rsplit(".", 1)
    if m := re.search(r"\d{4}-\d{2}-\d{2}", file):
        return d(m.group(0))
    # date = file.rsplit(".pdf", 1)[0]
    # if "-" in file:
    #     date = file.rsplit("-", 1).pop()
    # else:
    #     date = file.rsplit("_", 1).pop()
    if m := re.search(r"\d{6}", file):
        # thai date in briefing filenames
        date = m.group(0)
        return datetime.datetime(
            day=int(date[0:2]), month=int(date[2:4]), year=int(date[4:6]) - 43 + 2000
        )
    return None


def find_dates(content):
    # 7 - 13/11/2563
    dates = re.findall(r"([0-9]+)/([0-9]+)/([0-9]+)", content)
    dates = set(
        [
            datetime.datetime(day=int(d[0]), month=int(d[1]), year=int(d[2]) - 543)
            for d in dates
        ]
    )
    return sorted([d for d in dates])


def to_switching_date(dstr):
    "turning str 2021-01-02 into date but where m and d need to be switched"
    if not dstr:
        return None
    date = d(dstr).date()
    if date.day < 13 and date.month < 13:
        date = datetime.date(date.year, date.day, date.month)
    return date


def previous_date(end, day):
    "return a date before {end} by {day} days"
    start = end
    while start.day != int(day):
        start = start - datetime.timedelta(days=1)
    return start


def find_thai_date(content):
    "find thai date like '17 เม.ย. 2563' "
    m3 = re.search(r"([0-9]+) *([^ ]+) *(25[0-9][0-9])", content)
    d2, month, year = m3.groups()
    month = (
        THAI_ABBR_MONTHS.index(month) + 1
        if month in THAI_ABBR_MONTHS
        else THAI_FULL_MONTHS.index(month) + 1
        if month in THAI_FULL_MONTHS
        else None
    )
    date = datetime.datetime(year=int(year) - 543, month=month, day=int(d2))
    return date


def find_date_range(content):
    "Parse thai date ranges line '11-17 เม.ย. 2563' or '04/04/2563 12/06/2563'"
    m1 = re.search(
        r"([0-9]+)/([0-9]+)/([0-9]+) [-–] ([0-9]+)/([0-9]+)/([0-9]+)", content
    )
    m2 = re.search(r"([0-9]+) *[-–] *([0-9]+)/([0-9]+)/(25[0-9][0-9])", content)
    m3 = re.search(r"([0-9]+) *[-–] *([0-9]+) *([^ ]+) *(25[0-9][0-9])", content)
    if m1:
        d1, m1, y1, d2, m2, y2 = m1.groups()
        start = datetime.datetime(day=int(d1), month=int(m1), year=int(y1) - 543)
        end = datetime.datetime(day=int(d2), month=int(m2), year=int(y2) - 543)
        return start, end
    elif m2:
        d1, d2, month, year = m2.groups()
        end = datetime.datetime(year=int(year) - 543, month=int(month), day=int(d2))
        start = previous_date(end, d1)
        return start, end
    elif m3:
        d1, d2, month, year = m3.groups()
        month = (
            THAI_ABBR_MONTHS.index(month) + 1
            if month in THAI_ABBR_MONTHS
            else THAI_FULL_MONTHS.index(month) + 1
            if month in THAI_FULL_MONTHS
            else None
        )
        end = datetime.datetime(year=int(year) - 543, month=month, day=int(d2))
        start = previous_date(end, d1)
        return start, end
    else:
        return None, None


def daterange(start_date, end_date, offset=0):
    "return a range of dates from start_date until before end_date. Offset extends range by offset days"
    for n in range(int((end_date - start_date).days) + offset):
        yield start_date + datetime.timedelta(n)


def spread_date_range(start, end, row, columns):
    "take some values and spread it over a period of dates in proportion to data already there"
    r = list(daterange(start, end, offset=1))
    stats = [float(p) / len(r) for p in row]
    results = pd.DataFrame(
        [
            [
                date,
            ]
            + stats
            for date in r
        ],
        columns=columns,
    ).set_index("Date")
    return results


####################
# Extraction helpers
#####################


def parse_file(filename, html=False, paged=True):
    pages_txt = []

    # Read PDF file
    data = parser.from_file(filename, xmlContent=True)
    xhtml_data = BeautifulSoup(data["content"], features="lxml")
    if html and not paged:
        return xhtml_data
    pages = xhtml_data.find_all("div", attrs={"class": ["page", "slide-content"]})
    if not pages:
        if not paged:
            return repr(xhtml_data)
        else:
            return [repr(xhtml_data)]

    # TODO: slides are divided by slide-content and slide-master-content rather than being contained
    for i, content in enumerate(pages):
        # Parse PDF data using TIKA (xml/html)
        # It's faster and safer to create a new buffer than truncating it
        # https://stackoverflow.com/questions/4330812/how-do-i-clear-a-stringio-object
        _buffer = StringIO()
        _buffer.write(str(content))
        parsed_content = parser.from_buffer(_buffer.getvalue())
        if parsed_content["content"] is None:
            continue

        # Add pages
        text = parsed_content["content"].strip()
        if html:
            pages_txt.append(repr(content))
        else:
            pages_txt.append(text)
    if paged:
        return pages_txt
    else:
        return '\n\n\n'.join(pages_txt)


NUM_RE = re.compile(r"\d+(?:\,\d+)*(?:\.\d+)?")
INT_RE = re.compile(r"\d+(?:\,\d+)*")


def get_next_numbers(content, *matches, debug=False, before=False, remove=0, ints=True, until=None):
    if len(matches) == 0:
        matches = [""]
    for match in matches:
        if type(match) == str:
            match = re.compile(f"({match})")
        ahead, *behind = match.split(content, 1) if match else ("", "", content)
        if not behind:
            continue
        matched, *behind = behind
        behind = "".join(behind)
        found = ahead if before else behind
        if until is not None and until in found:
            found, rest = found.split(until, 1)  # TODO: how to put it back togeather if behind=True?
            rest = until+rest
        else:
            rest = ""
        numbers = (INT_RE if ints else NUM_RE).findall(found)
        numbers = [n.replace(",", "") for n in numbers]
        numbers = [int(n) if ints else float(n) for n in numbers if n]
        numbers = numbers if not before else list(reversed(numbers))
        if remove:
            behind = (INT_RE if ints else NUM_RE).sub("", found, remove)
        return numbers, matched + " " + rest + behind 
    if debug and matches:
        print("Couldn't find '{}'".format(match))
        print(content)
    return [], content


def get_next_number(content, *matches, default=None, remove=False, before=False, until=None):
    num, rest = get_next_numbers(content, *matches, remove=1 if remove else 0, before=before, until=until)
    return num[0] if num else default, rest


def slide2text(slide):
    text = ""
    if slide.shapes.title:
        text += slide.shapes.title.text
    for shape in slide.shapes:
        if shape.has_text_frame:
            # for p in shape.text_frame:
            text += "\n" + shape.text
    return text


def slide2chartdata(slide):
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart = shape.chart
        if chart is None:
            continue
        title = chart.chart_title.text_frame.text if chart.has_title else ""
        start, end = find_date_range(title)
        if start is None:
            continue
        series = dict([(s.name, s.values) for s in chart.series])

        yield chart, title, start, end, series


####################
# Download helpers
####################


def is_remote_newer(file, remote_date, check=True):
    if not os.path.exists(file):
        print(f"Missing: {file}")
        return True
    elif os.stat(file).st_size == 0:
        return True
    elif not check:
        return False
    elif remote_date is None:
        return True  # TODO: should we always keep cached?
    if type(remote_date) == str:
        remote_date = dateutil.parser.parse(remote_date).astimezone()
    fdate = datetime.datetime.fromtimestamp(os.path.getmtime(file)).astimezone()
    if remote_date > fdate:
        timestamp = fdate.strftime("%Y%m%d-%H%M%S")
        os.rename(file, f"{file}.{timestamp}")
        return True
    return False


def web_links(*index_urls, ext=".pdf", dir="html", match=None):
    def is_ext(a):
        return len(a.get("href").rsplit(ext)) == 2 if ext else True
    def is_match(a):
        return a.get("href") and is_ext(a) and (match.search(a.get_text(strip=True)) if match else True)
    for index_url in index_urls:
        for file, index in web_files(index_url, dir=dir, check=True):
            soup = parse_file(file, html=True, paged=False)
            links = (urllib.parse.urljoin(index_url, a.get('href')) for a in soup.find_all('a') if is_match(a))
            for l in links:
                yield l


def web_files(*urls, dir=os.getcwd(), check=CHECK_NEWER):
    "if check is None, then always download"
    for url in urls:
        modified = s.head(url).headers.get("last-modified") if check else None
        file = url.rsplit("/", 1)[-1]
        file = os.path.join(dir, file)
        os.makedirs(os.path.dirname(file), exist_ok=True)
        if is_remote_newer(file, modified, check):
            r = s.get(url)
            if r.status_code == 200:
                print(f"Download: {file}", end="")
                os.makedirs(os.path.dirname(file), exist_ok=True)
                with open(file, "wb") as f:
                    for chunk in r.iter_content(chunk_size=512 * 1024):
                        if chunk:  # filter out keep-alive new chunks
                            f.write(chunk)
                            print(".", end="")
                print("")
            elif os.path.exists(file):
                print(f"Error downloading: {file}: using cache")
            else:
                print(f"Error downloading: {file}: skipping")
                continue
        with open(file, "rb") as f:
            content = f.read()
        yield file, content


def dav_files(url="http://nextcloud.dmsc.moph.go.th/public.php/webdav", username="wbioWZAQfManokc", password="null",
              ext=".pdf .pptx", dir="testing_moph"):

    options = {
        "webdav_hostname": url,
        "webdav_login": username,
        "webdav_password": password,
    }
    client = Client(options)
    client.session.mount("http://", HTTPAdapter(max_retries=RETRY))
    client.session.mount("https://", HTTPAdapter(max_retries=RETRY))
    # important we get them sorted newest files first as we only fill in NaN from each additional file
    files = sorted(
        client.list(get_info=True),
        key=lambda info: dateutil.parser.parse(info["modified"]),
        reverse=True,
    )
    for info in files:
        file = info["path"].split("/")[-1]
        if not any([ext == file[-len(ext):] for ext in ext.split()]):
            continue
        target = os.path.join(dir, file)
        os.makedirs(os.path.dirname(target), exist_ok=True)
        if is_remote_newer(target, info["modified"]):
            client.download_file(file, target)
        yield target


##########################################
# download and parse thailand covid data
##########################################
d = dateutil.parser.parse


def situation_cases_cum(parsed_pdf, date):
    _, rest = get_next_numbers(parsed_pdf, "The Disease Situation in Thailand", debug=True)
    cases, rest = get_next_numbers(
        rest, 
        "Total number of confirmed cases",
        "Characteristics of Infection in Confirmed cases",
        "Confirmed cases",
        debug=False
    )
    if not cases:
        return pd.DataFrame()
    cases, *_ = cases
    if date < d("2020-04-09"):
        return pd.DataFrame([(date, cases)], columns=["Date", "Cases Cum"]).set_index("Date")
    outside_quarantine, _ = get_next_numbers(
        rest,
        "Cases found outside (?:the )?(?:state )?quarantine (?:facilities|centers)",
        debug=False
    )
    if outside_quarantine:
        outside_quarantine, *_ = outside_quarantine
        # 2647.0 # 2021-02-15
        # if date > d("2021-01-25"):
        #    # thai graphic says imported is 2396 instead of en 4195 on 2021-01-26
        #    outside_quarantine = outside_quarantine -  (4195 - 2396) 

        quarantine, _ = get_next_number(
            rest, 
            "Cases found in (?:the )?(?:state )?quarantine (?:facilities|centers)",
            "Staying in [^ ]* quarantine",
            default=0, until="●")
        quarantine = 1903 if quarantine == 19003 else quarantine  # "2021-02-05"
        # TODO: work out date when it flips back again.
        if date == d("2021-05-17"):
            imported = quarantine = outside_quarantine
            outside_quarantine = 0
        elif date < d("2020-12-28") or (date > d("2021-01-25") and outside_quarantine > quarantine):
            imported = outside_quarantine  # It's mislabeled (new daily is correct however)
            imported = 2647 if imported == 609 else imported  # "2021-02-17")
            imported = None if imported == 610 else imported  # 2021-02-20 - 2021-03-01
            if imported is not None: 
                outside_quarantine = imported - quarantine
            else:
                outside_quarantine = None
        else:
            imported = outside_quarantine + quarantine
    else:
        quarantine, _ = get_next_numbers(
            rest, 
            "(?i)d.?e.?signated quarantine",
            debug=False)
        quarantine, *_ = quarantine if quarantine else [None]
        quarantine = 562 if quarantine == 5562 else quarantine  # "2021-09-19"
        imported, _ = get_next_number(
            rest, 
            "(?i)Imported Case(?:s)?",
            "(?i)Cases were imported from overseas")
        if imported and quarantine:
            outside_quarantine = imported - quarantine
        else:
            outside_quarantine = None  #TODO: can we get imported from total - quarantine - local?
    if quarantine:
        active, _ = get_next_number(
            rest,
            "(?i)Cases found from active case finding",
            "(?i)Cases were (?:infected )?migrant workers",
        )
        prison, _ = get_next_number(rest, "Cases found in Prisons", default=0)
        if active is not None:
            active += prison

        # TODO: cum local really means all local ie walkins+active testing 
        local, _ = get_next_number(rest, "(?i)(?:Local )?Transmission")
        # TODO: 2021-01-25. Local 6629.0 -> 12250.0, quarantine 597.0 -> 2396.0 active 4684.0->5532.0
        if imported is None:
            pass
        elif cases-imported == active:
            walkin = local
            local = cases-imported
            active = local - walkin
        elif active is None:
            pass
        elif local + active == cases-imported:
            # switched to different definition?
            walkin = local
            local = walkin + active
        elif date <= d("2021-01-25") or d("2021-02-16") <= date <= d("2021-03-01"):
            walkin = local
            local = walkin + active
    else:
        local, active = None, None
    
    #assert cases == (local+imported) # Too many mistakes
    return pd.DataFrame(
        [(date, cases, local, imported, quarantine, outside_quarantine, active)],
        columns=["Date", "Cases Cum", "Cases Local Transmission Cum", "Cases Imported Cum", "Cases In Quarantine Cum", "Cases Outside Quarantine Cum", "Cases Proactive Cum"]
        ).set_index("Date")


def situation_cases_new(parsed_pdf, date):
    if date < d("2020-11-02"):
        return pd.DataFrame()
    _, rest = get_next_numbers(
        parsed_pdf,
        "The Disease Situation in Thailand", 
        "(?i)Type of case Total number Rate of Increase",
        debug=False)
    cases, rest = get_next_numbers(
        rest, 
        "(?i)number of new case(?:s)?",
        debug=False
    )
    if not cases or date < d("2020-05-09"):
        return pd.DataFrame()
    cases, *_ = cases
    local, _ = get_next_numbers(rest, "(?i)(?:Local )?Transmission", debug=False)
    local, *_ = local if local else [None]
    quarantine, _ = get_next_numbers(
        rest, 
        "Cases found (?:positive from |in )(?:the )?(?:state )?quarantine",
        #"Staying in [^ ]* quarantine",
        debug=False)
    quarantine, *_ = quarantine
    quarantine = {d("2021-01-27"): 11}.get(date, quarantine)  # corrections from thai doc
    outside_quarantine, _ = get_next_numbers(
        rest,
        "(?i)Cases found (?:positive )?outside (?:the )?(?:state )?quarantine",
        debug=False
    )
    outside_quarantine, *_ = outside_quarantine if outside_quarantine else [None]
    if outside_quarantine is not None:
        imported = quarantine + outside_quarantine
        active, _ = get_next_numbers(
            rest,
            "(?i)active case",
            debug=True
        )
        active, *_ = active if active else [None]
        if date <= d("2020-12-24"):  # starts getting cum values
            active = None
        # local really means walkins. so need add it up
        if active:
            local = local + active
    else:
        imported, active = None, None
    cases = {d("2021-03-31"): 42}.get(date, cases)
    #if date not in [d("2020-12-26")]:
    #    assert cases == (local+imported) # except 2020-12-26 - they didn't include 30 proactive
    return pd.DataFrame(
        [(date, cases, local, imported, quarantine, outside_quarantine, active)],
        columns=["Date", "Cases", "Cases Local Transmission", "Cases Imported", "Cases In Quarantine", "Cases Outside Quarantine", "Cases Proactive"]
        ).set_index("Date")


re_walkin_priv = re.compile("(?i)cases (in|at) private hospitals")


def situation_pui(parsed_pdf, date):
    numbers, _ = get_next_numbers(
        parsed_pdf, "Total +number of laboratory tests",
        until="Sought medical services on their own at hospitals", 
        debug=False
    )
    if numbers:
        if len(numbers) == 7:
            tests_total, pui, active_finding, asq, not_pui, pui2, pui_port, *rest = numbers
        elif len(numbers) == 6:
            tests_total, pui, asq, not_pui, pui2, pui_port, *rest = numbers
            active_finding = None
        else:
            raise Exception(numbers)

        pui = {309371: 313813}.get(pui, pui)  # 2020-07-01
        #TODO: find 1529045 below and see which is correct 20201-04-26
        pui2 = pui if pui2 in [96989, 433807, 3891136, 385860, 326073, 1529045] else pui2
        assert pui == pui2
    else:
        numbers, _ = get_next_numbers(
            parsed_pdf, "Total number of people who met the criteria of patients", debug=False,
        )
        if date > dateutil.parser.parse("2020-01-30") and not numbers:
            raise Exception(f"Problem parsing {date}")
        elif not numbers:
            return pd.DataFrame()
        tests_total, active_finding, asq, not_pui = [None] * 4
        pui, pui_airport, pui_seaport, pui_hospital, *rest = numbers
        # pui_port = pui_airport + pui_seaport
    if pui in [1103858, 3891136]:  # mistypes? # 433807?
        pui = None
    elif tests_total in [783679, 849874, 936458]:
        tests_total = None
    elif None in (tests_total, pui, asq, not_pui) and date > d("2020-06-30"):
        raise Exception(f"Missing data at {date}")

    # walkin public vs private
    numbers, rest = get_next_numbers(parsed_pdf, "Sought medical services on their own at hospitals")
    if not numbers:
        pui_walkin_private, pui_walkin_public, pui_walkin = [None]*3
    elif re_walkin_priv.search(rest):
        pui_walkin_private, pui_walkin_public, pui_walkin, *_ = numbers
        pui_walkin_public = {8628765: 862876}.get(pui_walkin_public, pui_walkin_public)
        #assert pui_walkin == pui_walkin_private + pui_walkin_public
    else:
        pui_walkin, *_ = numbers
        pui_walkin_private, pui_walkin_public = None, None
        pui_walkin = {853189: 85191}.get(pui_walkin, pui_walkin)  # by taking away other numbers
    assert pui_walkin is None or pui is None or (pui_walkin <= pui and 5000000 > pui_walkin > 0)
    assert pui_walkin_public is None or (5000000 > pui_walkin_public > 10000)

    row = (tests_total, pui, active_finding, asq, not_pui, pui_walkin, pui_walkin_private, pui_walkin_public)
    return pd.DataFrame(
        [(date, )+row],
        columns=[
            "Date", 
            "Tested Cum", 
            "Tested PUI Cum", 
            "Tested Proactive Cum", 
            "Tested Quarantine Cum", 
            "Tested Not PUI Cum",
            "Tested PUI Walkin Cum",
            "Tested PUI Walkin Private Cum",
            "Tested PUI Walkin Public Cum",
            ]
        ).set_index("Date")


def get_en_situation():
    results = pd.DataFrame(columns=["Date"]).set_index("Date")
    url = "https://ddc.moph.go.th/viralpneumonia/eng/situation.php"
    for file, _ in web_files(*web_links(url, ext=".pdf", dir="situation_en"), dir="situation_en"):
        parsed_pdf = parse_file(file, html=False, paged=False).replace("\u200b", "")
        if "situation" not in os.path.basename(file):
            continue
        date = file2date(file)
        if date <= dateutil.parser.parse("2020-01-30"):
            continue  # TODO: can manually put in numbers before this
        pui = situation_pui(parsed_pdf, date)
        cases = situation_cases_cum(parsed_pdf, date)
        new_cases = situation_cases_new(parsed_pdf, date)
        row = pui.combine_first(cases).combine_first(new_cases)
        results = results.combine_first(row)
        # cums = [c for c in results.columns if ' Cum' in c]
        # if len(results) > 1 and (results.iloc[0][cums] > results.iloc[1][cums]).any():
        #     print((results.iloc[0][cums] > results.iloc[1][cums]))
        #     print(results.iloc[0:2])
        # raise Exception("Cumulative data didn't increase")
        # row = results.iloc[0].to_dict()
        print(date.date(), file, row.to_string(header=False, index=False))
        #     "p{Tested PUI Cum:.0f}\tc{Cases Cum:.0f}({Cases:.0f})\t"
        #     "l{Cases Local Transmission Cum:.0f}({Cases Local Transmission:.0f})\t"
        #     "a{Cases Proactive Cum:.0f}({Cases Proactive:.0f})\t"
        #     "i{Cases Imported Cum:.0f}({Cases Imported:.0f})\t"
        #     "q{Cases In Quarantine Cum:.0f}({Cases In Quarantine:.0f})\t"
        #     "".format(**row)
        # )
    # Missing data. filled in from th infographic
    missing = [
        (d("2020-12-19"), 2476, 0, 0),
        (d("2020-12-20"), 3011, 516, 516),
        (d("2020-12-21"), 3385, 876, 360),
        (d("2020-12-22"), 3798, 1273, 397),
        (d("2020-12-23"), 3837, 1273, 0),
        (d("2020-12-24"), 3895, 1273, 0),
        (d("2020-12-25"), 3976, 1308, 35),
    ]
    missing = pd.DataFrame(
        missing,
        columns=["Date", "Cases Local Transmission Cum", "Cases Proactive Cum", "Cases Proactive"]
    ).set_index("Date")
    results = missing[["Cases Local Transmission Cum", "Cases Proactive Cum", ]].combine_first(results)
    return results


def get_situation_today():
    _, page = next(web_files("https://ddc.moph.go.th/viralpneumonia/index.php", dir="situation_th", check=True))
    text = BeautifulSoup(page, 'html.parser').get_text()
    numbers, rest = get_next_numbers(text, "ผู้ป่วยเข้าเกณฑ์เฝ้าระวัง")
    pui_cum, pui = numbers[:2]
    numbers, rest = get_next_numbers(text, "กักกันในพื้นที่ที่รัฐกำหนด")
    imported_cum, imported = numbers[:2]
    numbers, rest = get_next_numbers(text, "ผู้ป่วยยืนยัน")
    cases_cum, cases = numbers[:2]
    numbers, rest = get_next_numbers(text, "สถานการณ์ในประเทศไทย")
    date = find_thai_date(rest).date()
    row = [cases_cum, cases, pui_cum, pui, imported_cum, imported]
    return pd.DataFrame(
        [[date, ]+row],
        columns=["Date", "Cases Cum", "Cases", "Tested PUI Cum", "Tested PUI", "Cases Imported Cum", "Cases Imported"]
    ).set_index("Date")
    

def check_cum(df, results):
    if results.empty:
        return True
    next_day = results.loc[results.index[0]][[c for c in results.columns if " Cum" in c]]
    last = df.loc[df.index[-1]][[c for c in df.columns if " Cum" in c]]
    if (next_day.fillna(0) >= last.fillna(0)).all():
        return True
    else:
        raise Exception(str(next_day - last))


def situation_pui_th(parsed_pdf, date, results):
    tests_total, active_finding, asq, not_pui = [None] * 4
    numbers, content = get_next_numbers(
        parsed_pdf,
        r"ด่านโรคติดต่อระหว่างประเทศ",
        r"ด่านโรคติดต่อระหวา่งประเทศ",  # 'situation-no346-141263n.pdf'
        r"นวนการตรวจทาง\S+องปฏิบัติการ",
        "ด่านควบคุมโรคติดต่อระหว่างประเทศ",
    )
    # cases = None
    if len(numbers) > 6:
        (
            screened_port,
            screened_cw,
            tests_total,
            pui,
            active_finding,
            asq,
            not_pui,
            *rest,
        ) = numbers
        if tests_total < 30000:
            tests_total, pui, active_finding, asq, not_pui, *rest = numbers
            if pui == 4534137:
                pui = 453413  # situation-no273-021063n.pdf
    else:
        numbers, content = get_next_numbers(
            parsed_pdf,
            # "ผู้ป่วยที่มีอาการเข้าได้ตามนิยาม",
            "ตาราง 2 ผลดำ",
            "ตาราง 1",  # situation-no172-230663.pdf #'situation-no83-260363_1.pdf'
        )
        if len(numbers) > 0:
            pui, *rest = numbers
    if date > dateutil.parser.parse("2020-03-26") and not numbers:
        raise Exception(f"Problem parsing {date}")
    elif not numbers:
        return
    if tests_total == 167515:  # situation-no447-250364.pdf
        tests_total = 1675125
    if date in [d("2020-12-23")]:  # 1024567
        tests_total, not_pui = 997567, 329900
    if (
        tests_total is not None
        and tests_total > 2000000 < 30000
        or pui > 1500000 < 100000
    ):
        raise Exception(f"Bad data in {date}")
    pui = {d("2020-02-12"): 799, d("2020-02-13"): 804}.get(date, pui)  # Guess

    walkinsre = "(?:ษาที่โรงพยาบาลด้วยตนเอง|โรงพยาบาลด้วยตนเอง|ารับการรักษาท่ีโรงพยาบาลด|โรงพยาบาลดวยตนเอง)"
    _, line = get_next_numbers(parsed_pdf, walkinsre)
    pui_walkin_private, rest = get_next_number(line, f"(?s){walkinsre}.*?โรงพยาบาลเอกชน", remove=True)
    pui_walkin_public, rest = get_next_number(rest, f"(?s){walkinsre}.*?โรงพยาบาลรัฐ", remove=True)
    unknown, rest = get_next_number(rest, f"(?s){walkinsre}.*?(?:งการสอบสวน|ารสอบสวน)", remove=True)
    #rest = re.sub("(?s)(?:งการสอบสวน|ารสอบสวน).*?(?:อ่ืนๆ|อื่นๆ|อืน่ๆ|ผู้ป่วยยืนยันสะสม|88)?", "", rest,1)
    pui_walkin, rest = get_next_number(rest)
    assert pui_walkin is not None
    if date <= d("2020-03-10"):
        pui_walkin_private, pui_walkin, pui_walkin_public = [None] * 3  # starts going up again
    #pui_walkin_private = {d("2020-03-10"):2088}.get(date, pui_walkin_private)

    assert pui_walkin is None or pui is None or (pui_walkin <= pui and pui_walkin > 0)

    row = (tests_total, pui, active_finding, asq, not_pui, pui_walkin_private, pui_walkin_public, pui_walkin)
    if None in row and date > d("2020-06-30"):
        raise Exception(f"Missing data at {date}")
    df = pd.DataFrame(
         [(date,)+row],
         columns=[
             "Date",
             "Tested Cum",
             "Tested PUI Cum",
             "Tested Proactive Cum",
             "Tested Quarantine Cum",
             "Tested Not PUI Cum",
             "Tested PUI Walkin Private Cum",
             "Tested PUI Walkin Public Cum",
             "Tested PUI Walkin Cum"]
    ).set_index("Date")
    assert check_cum(df, results)
    return df
     

def get_thai_situation():
    results = pd.DataFrame(columns=["Date"]).set_index("Date")
    links = web_links(        
        "https://ddc.moph.go.th/viralpneumonia/situation.php",
        "https://ddc.moph.go.th/viralpneumonia/situation_more.php",
        ext=".pdf",
        dir="situation_th"
    )
    for file, _ in web_files(*links, dir="situation_th"):
        parsed_pdf = parse_file(file, html=False, paged=False)
        if "situation" not in os.path.basename(file):
            continue
        if "Situation Total number of PUI" in parsed_pdf:
            # english report mixed up? - situation-no171-220663.pdf
            continue
        date = file2date(file)
        df = situation_pui_th(parsed_pdf, date, results)
        if df is not None:
            results = results.combine_first(df)
            print(date.date(), file, df.to_string(header=False, index=False))
            # file,
            # "p{Tested PUI Cum:.0f}\t"
            # "t{Tested Cum:.0f}\t"
            # "{Tested Proactive Cum:.0f}\t"
            # "{Tested Quarantine Cum:.0f}\t"
            # "{Tested Not PUI Cum:.0f}\t"
            # "".format(**results.iloc[0].to_dict()))

    #print(results)
    return results


def cum2daily(results):
    cum = results[(c for c in results.columns if " Cum" in c)]
    all_days = pd.date_range(cum.index.min(), cum.index.max(), name="Date")
    cum = cum.reindex(all_days)  # put in missing days with NaN
    #cum = cum.interpolate(limit_area="inside") # missing dates need to be filled so we don't get jumps
    cum = cum - cum.shift(+1)  # we got cumilitive data
    renames = dict((c, c.rstrip(' Cum')) for c in list(cum.columns) if 'Cum' in c)
    cum = cum.rename(columns=renames)
    return cum


def get_situation():
    print("========Situation Reports==========")

    today_situation = get_situation_today()
    en_situation = get_en_situation()
    th_situation = get_thai_situation()
    situation = th_situation.combine_first(en_situation)
    cum = cum2daily(situation)
    situation = situation.combine_first(cum)  # any direct non-cum are trusted more

    # Only add in the live stats if they have been updated with new info
    today = today_situation.index.max()
    yesterday = today-datetime.timedelta(days=1)
    stoday = today_situation.loc[today]
    syesterday = situation.loc[str(yesterday)] if str(yesterday) in situation else None
    if syesterday is None or (syesterday['Tested PUI Cum'] < stoday['Tested PUI Cum'] and syesterday['Tested PUI'] != stoday['Tested PUI']):
        situation = situation.combine_first(today_situation)

    export(situation, "situation_reports")
    return situation


def get_cases():
    print("========Covid19 Timeline==========")

    file, text = next(web_files("https://covid19.th-stat.com/api/open/timeline?123", dir="json", check=True))
    data = pd.DataFrame(json.loads(text)['Data'])
    data['Date'] = pd.to_datetime(data['Date'])
    data = data.set_index("Date")
    cases = data[["NewConfirmed", "NewDeaths", "NewRecovered", "Hospitalized"]]
    cases = cases.rename(columns=dict(NewConfirmed="Cases", NewDeaths="Deaths", NewRecovered="Recovered"))
    return cases


def get_tests_by_day():
    print("========Tests by Day==========")

    file = next(dav_files(ext="xlsx"))
    tests = pd.read_excel(file, parse_dates=True, usecols=[0, 1, 2])
    tests.dropna(how="any", inplace=True)  # get rid of totals row
    tests = tests.set_index("Date")
    pos = tests.loc["Cannot specify date"].Pos
    total = tests.loc["Cannot specify date"].Total
    tests.drop("Cannot specify date", inplace=True)
    # Need to redistribute the unknown values across known values
    # Documentation tells us it was 11 labs and only before 3 April
    unknown_end_date = datetime.datetime(day=3, month=4, year=2020)
    all_pos = tests["Pos"][:unknown_end_date].sum()
    all_total = tests["Total"][:unknown_end_date].sum()
    for index, row in tests.iterrows():
        if index > unknown_end_date:
            continue
        row.Pos = float(row.Pos) + row.Pos / all_pos * pos
        row.Total = float(row.Total) + row.Total / all_total * total
    # TODO: still doesn't redistribute all missing values due to rounding. about 200 left
    #print(tests["Pos"].sum(), pos + all_pos)
    #print(tests["Total"].sum(), total + all_total)
    # fix datetime
    tests.reset_index(drop=False, inplace=True)
    tests["Date"] = pd.to_datetime(tests["Date"])
    tests.set_index("Date", inplace=True)

    tests.rename(columns=dict(Pos="Pos XLS", Total="Tests XLS"), inplace=True)
    print(file, len(tests))

    return tests


DISTRICT_RANGE_SIMPLE = [str(i) for i in range(1, 14)]
DISTRICT_RANGE = DISTRICT_RANGE_SIMPLE + ["Prison"]
DISTRICT_RANGE_UNKNOWN = [str(i) for i in range(1, 14)] + ["Prison", "Unknown"]


def get_tests_by_area():
    pos_cols = [f"Pos Area {i}" for i in DISTRICT_RANGE_SIMPLE]
    test_cols = [f"Tests Area {i}" for i in DISTRICT_RANGE_SIMPLE]
    columns = ["Date"] + pos_cols + test_cols + ["Pos Area", "Tests Area"]
    raw_cols = ["Start", "End", ] + pos_cols + test_cols
    data = pd.DataFrame()
    raw = pd.DataFrame()

    for file in dav_files(ext=".pptx"):
        prs = Presentation(file)
        for chart in (
            chart for slide in prs.slides for chart in slide2chartdata(slide)
        ):
            chart, title, start, end, series = chart
            if "เริ่มเปิดบริการ" not in title and any(
                t in title for t in ["เขตสุขภาพ", "เขตสุขภำพ"]
            ):
                # the graph for X period split by health area.
                # Need both pptx and pdf as one pdf is missing
                pos = list(series["จำนวนผลบวก"])
                tests = list(series["จำนวนตรวจ"])
                row = pos + tests + [sum(pos), sum(tests)]
                results = spread_date_range(start, end, row, columns)
                # print(results)
                data = data.combine_first(results)
                raw = raw.combine_first(pd.DataFrame(
                    [[start, end, ]+pos + tests],
                    columns=raw_cols
                ).set_index("Start"))
        print(file)
    # Also need pdf copies becaus of missing pptx
    for file in dav_files(ext=".pdf"):
        pages = parse_file(file, html=False, paged=True)
        not_whole_year = [page for page in pages if "เริ่มเปิดบริการ" not in page]
        by_area = [
            page
            for page in not_whole_year
            if "เขตสุขภาพ" in page or "เขตสุขภำพ" in page
        ]
        # Can't parse '35_21_12_2020_COVID19_(ถึง_18_ธันวาคม_2563)(powerpoint).pptx' because data is a graph
        # no pdf available so data missing
        # Also missing 14-20 Nov 2020 (no pptx or pdf)

        for page in by_area:
            start, end = find_date_range(page)
            if start is None:
                continue
            if "349585" in page:
                page = page.replace("349585", "349 585")
            # if '16/10/2563' in page:
            #     print(page)
            # First line can be like จดัท ำโดย เพญ็พชิชำ ถำวงศ ์กรมวิทยำศำสตณก์ำรแพทย์ วันที่ท ำรำยงำน 15/02/2564 เวลำ 09.30 น.
            first, rest = page.split("\n", 1)
            page = (
                rest if "เพญ็พชิชำ" in first or "/" in first else page
            )  # get rid of first line that sometimes as date and time in it
            numbers, content = get_next_numbers(page, "", debug=True)  # "ภาคเอกชน",
            # ภาครัฐ
            # ภาคเอกชน
            # จดัท ำโดย เพญ็พชิชำ ถำวงศ ์กรมวิทยำศำสตณก์ำรแพทย์
            # print(numbers)
            # TODO: should really find and parse X axis labels which contains 'เขต' and count
            tests_start = 13 if "total" not in page else 14
            pos = numbers[0:13]
            tests = numbers[tests_start:tests_start + 13]
            row = pos + tests + [sum(pos), sum(tests)]
            results = spread_date_range(start, end, row, columns)
            #print(results)
            data = data.combine_first(results)
            raw = raw.combine_first(pd.DataFrame(
                [[start, end, ]+pos + tests],
                columns=raw_cols
            ).set_index("Start"))
        print(file)
    export(raw, "tests_by_area")
    return data


def get_tests_private_public():
    print("========Tests public+private==========")

    data = pd.DataFrame()

    # some additional data from pptx files
    for file in dav_files(ext=".pptx"):
        prs = Presentation(file)
        for chart in (
            chart for slide in prs.slides for chart in slide2chartdata(slide)
        ):
            chart, title, start, end, series = chart
            if "เริ่มเปิดบริการ" not in title and any(
                t in title for t in ["เขตสุขภาพ", "เขตสุขภำพ"]
            ):
                # area graph
                continue
            elif "และอัตราการตรวจพบ" in title and "รายสัปดาห์" not in title:
                # The graphs at the end with all testing numbers private vs public
                private = " Private" if "ภาคเอกชน" in title else ""

                # pos = series["Pos"]
                if "จำนวนตรวจ" not in series:
                    continue
                tests = series["จำนวนตรวจ"]
                positivity = series["% Detection"]
                dates = list(daterange(start, end, 1))
                df = pd.DataFrame(
                    {
                        "Date": dates,
                        f"Tests{private}": tests,
                        f"% Detection{private}": positivity,
                    }
                ).set_index("Date")
                df[f"Pos{private}"] = (
                    df[f"Tests{private}"] * df[f"% Detection{private}"] / 100.0
                )
                #print(df)
                data = data.combine_first(df)
            # TODO: There is also graphs splt by hospital
        print(file, len(data))
        if not data.empty:
            break  # we only need the latest
    data['Pos Public'] = data['Pos'] - data['Pos Private']
    data['Tests Public'] = data['Tests'] - data['Tests Private']
    export(data, "tests_pubpriv")
    return data


def get_provinces():
    #_, districts = next(web_files("https://en.wikipedia.org/wiki/Healthcare_in_Thailand#Health_Districts", dir="html"))
    areas = pd.read_html("https://en.wikipedia.org/wiki/Healthcare_in_Thailand#Health_Districts")[0]
    provinces = areas.assign(Provinces=areas['Provinces'].str.split(", ")).explode("Provinces")
    provinces['Provinces'] = provinces['Provinces'].str.strip()
    provinces = provinces.rename(columns=dict(Provinces="ProvinceEn")).drop(columns="Area Code")
    provinces['ProvinceAlt'] = provinces['ProvinceEn']
    provinces = provinces.set_index("ProvinceAlt")
    provinces.loc["Bangkok"] = [13, "Central", "Bangkok"]
    provinces.loc["Unknown"] = ["Unknown", "", "Unknown"]
    provinces.loc["Prison"] = ["Prison", "", "Prison"]
    provinces['Health District Number'] = provinces['Health District Number'].astype(str)

    # extra spellings for matching
    provinces.loc['Korat'] = provinces.loc['Nakhon Ratchasima']
    provinces.loc['Khorat'] = provinces.loc['Nakhon Ratchasima']
    provinces.loc['Suphanburi'] = provinces.loc['Suphan Buri']
    provinces.loc["Ayutthaya"] = provinces.loc["Phra Nakhon Si Ayutthaya"]
    provinces.loc["Phathum Thani"] = provinces.loc["Pathum Thani"]
    provinces.loc["Pathom Thani"] = provinces.loc["Pathum Thani"]
    provinces.loc["Ubon Thani"] = provinces.loc["Udon Thani"]
    provinces.loc["Bung Kan"] = provinces.loc["Bueng Kan"]
    provinces.loc["Chainat"] = provinces.loc["Chai Nat"]
    provinces.loc["Chon Buri"] = provinces.loc["Chonburi"]
    provinces.loc["ลาปาง"] = provinces.loc["Lampang"]
    provinces.loc["หนองบัวลาภู"] = provinces.loc["Nong Bua Lamphu"]
    provinces.loc["ปทุุมธานี"] = provinces.loc["Pathum Thani"]
    provinces.loc["เพชรบุรีี"] = provinces.loc["Phetchaburi"]
    provinces.loc["เพชรบุรีี"] = provinces.loc["Phetchaburi"]
    provinces.loc["เพชรบุุรี"] = provinces.loc["Phetchaburi"]

    provinces.loc["สมุุทรสาคร"] = provinces.loc["Samut Sakhon"]
    provinces.loc["สมุทธสาคร"] = provinces.loc["Samut Sakhon"]
    provinces.loc["กรุงเทพฯ"] = provinces.loc["Bangkok"]
    provinces.loc["กรุงเทพ"] = provinces.loc["Bangkok"]
    provinces.loc["กรงุเทพมหานคร"] = provinces.loc["Bangkok"]
    provinces.loc["พระนครศรีอยุธา"] = provinces.loc["Ayutthaya"]
    provinces.loc["อยุธยา"] = provinces.loc["Ayutthaya"]
    provinces.loc["สมุุทรสงคราม"] = provinces.loc["Samut Songkhram"]
    provinces.loc["สมุุทรปราการ"] = provinces.loc["Samut Prakan"]
    provinces.loc["สระบุุรี"] = provinces.loc["Saraburi"]
    provinces.loc["พม่า"] = provinces.loc["Nong Khai"]  # it's really burma, but have to put it somewhere
    provinces.loc["ชลบุุรี"] = provinces.loc["Chon Buri"]
    provinces.loc["นนทบุุรี"] = provinces.loc["Nonthaburi"]
    # from prov table in briefings
    provinces.loc["เชยีงใหม่"] = provinces.loc["Chiang Mai"]
    provinces.loc['จนัทบรุ'] = provinces.loc['Chanthaburi']
    provinces.loc['บรุรีมัย'] = provinces.loc['Buriram']
    provinces.loc['กาญจนบรุ'] = provinces.loc['Kanchanaburi']
    provinces.loc['Prachin Buri'] = provinces.loc['Prachinburi']
    provinces.loc['ปราจนีบรุ'] = provinces.loc['Prachin Buri']
    provinces.loc['ลพบรุ'] = provinces.loc['Lopburi']
    provinces.loc['พษิณุโลก'] = provinces.loc['Phitsanulok']
    provinces.loc['นครศรธีรรมราช'] = provinces.loc['Nakhon Si Thammarat']
    provinces.loc['เพชรบรูณ์'] = provinces.loc['Phetchabun']
    provinces.loc['อา่งทอง'] = provinces.loc['Ang Thong']
    provinces.loc['ชยัภมู'] = provinces.loc['Chaiyaphum']
    provinces.loc['รอ้ยเอ็ด'] = provinces.loc['Roi Et']
    provinces.loc['อทุยัธานี'] = provinces.loc['Uthai Thani']
    provinces.loc['ชลบรุ'] = provinces.loc['Chon Buri']
    provinces.loc['สมทุรปราการ'] = provinces.loc['Samut Prakan']
    provinces.loc['นราธวิาส'] = provinces.loc['Narathiwat']
    provinces.loc['ประจวบครีขีนัธ'] = provinces.loc['Prachuap Khiri Khan']
    provinces.loc['สมทุรสาคร'] = provinces.loc['Samut Sakhon']
    provinces.loc['ปทมุธานี'] = provinces.loc['Pathum Thani']
    provinces.loc['สระแกว้'] = provinces.loc['Sa Kaeo']
    provinces.loc['นครราชสมีา'] = provinces.loc['Nakhon Ratchasima']
    provinces.loc['นนทบรุ'] = provinces.loc['Nonthaburi']
    provinces.loc['ภเูก็ต'] = provinces.loc['Phuket']
    provinces.loc['สพุรรณบรุี'] = provinces.loc['Suphan Buri']
    provinces.loc['อดุรธานี'] = provinces.loc['Udon Thani']
    provinces.loc['พระนครศรอียธุยา'] = provinces.loc['Ayutthaya']
    provinces.loc['สระบรุ'] = provinces.loc['Saraburi']
    provinces.loc['เพชรบรุ'] = provinces.loc['Phetchaburi']
    provinces.loc['ราชบรุ'] = provinces.loc['Ratchaburi']
    provinces.loc['เชยีงราย'] = provinces.loc['Chiang Rai']
    provinces.loc['อบุลราชธานี'] = provinces.loc['Ubon Ratchathani']
    provinces.loc['สรุาษฎรธ์านี'] = provinces.loc['Surat Thani']
    provinces.loc['ฉะเชงิเทรา'] = provinces.loc['Chachoengsao']
    provinces.loc['สมทุรสงคราม'] = provinces.loc['Samut Songkhram']
    provinces.loc['แมฮ่อ่งสอน'] = provinces.loc['Mae Hong Son']
    provinces.loc['สโุขทยั'] = provinces.loc['Sukhothai']
    provinces.loc['นา่น'] = provinces.loc['Nan']
    provinces.loc['อตุรดติถ'] = provinces.loc['Uttaradit']
    provinces.loc['Nong Bua Lam Phu'] = provinces.loc['Nong Bua Lamphu']
    provinces.loc['หนองบวัล'] = provinces.loc['Nong Bua Lam Phu']
    provinces.loc['พงังา'] = provinces.loc['Phang Nga']
    provinces.loc['สรุนิทร'] = provinces.loc['Surin']
    provinces.loc['Si Sa Ket'] = provinces.loc['Sisaket']
    provinces.loc['ศรสีะเกษ'] = provinces.loc['Si Sa Ket']
    provinces.loc['ตรงั'] = provinces.loc['Trang']
    provinces.loc['พจิติร'] = provinces.loc['Phichit']
    provinces.loc['ปตัตานี'] = provinces.loc['Pattani']
    provinces.loc['ชยันาท'] = provinces.loc['Chai Nat']
    provinces.loc['พทัลงุ'] = provinces.loc['Phatthalung']
    provinces.loc['มกุดาหาร'] = provinces.loc['Mukdahan']
    provinces.loc['บงึกาฬ'] = provinces.loc['Bueng Kan']
    provinces.loc['กาฬสนิธุ'] = provinces.loc['Kalasin']
    provinces.loc['สงิหบ์รุ'] = provinces.loc['Sing Buri']
    provinces.loc['ปทมุธาน'] = provinces.loc['Pathum Thani']
    provinces.loc['สพุรรณบรุ'] = provinces.loc['Suphan Buri']
    provinces.loc['อดุรธาน'] = provinces.loc['Udon Thani']
    provinces.loc['อบุลราชธาน'] = provinces.loc['Ubon Ratchathani']
    provinces.loc['สรุาษฎรธ์าน'] = provinces.loc['Surat Thani']
    provinces.loc['นครสวรรค'] = provinces.loc['Nakhon Sawan']
    provinces.loc['ล าพนู'] = provinces.loc['Lamphun']
    provinces.loc['ล าปาง'] = provinces.loc['Lampang']
    provinces.loc['เพชรบรูณ'] = provinces.loc['Phetchabun']
    provinces.loc['อทุยัธาน'] = provinces.loc['Uthai Thani']
    provinces.loc['ก าแพงเพชร'] = provinces.loc['Kamphaeng Phet']
    provinces.loc['Lam Phu'] = provinces.loc['Nong Bua Lamphu']
    provinces.loc['หนองบวัล าภู'] = provinces.loc['Lam Phu']
    provinces.loc['ปตัตาน'] = provinces.loc['Pattani']
    provinces.loc['บรุรีัมย'] = provinces.loc['Buriram']
    provinces.loc['Buri Ram'] = provinces.loc['Buriram']
    provinces.loc['บรุรัีมย'] = provinces.loc['Buri Ram']
    provinces.loc['จันทบรุ'] = provinces.loc['Chanthaburi']
    provinces.loc['ชมุพร'] = provinces.loc['Chumphon']
    provinces.loc['อทุัยธาน'] = provinces.loc['Uthai Thani']
    provinces.loc['อ านาจเจรญ'] = provinces.loc['Amnat Charoen']
    provinces.loc['สโุขทัย'] = provinces.loc['Sukhothai']
    provinces.loc['ปัตตาน'] = provinces.loc['Pattani']
    provinces.loc['พัทลงุ'] = provinces.loc['Phatthalung']
    provinces.loc['ขอนแกน่'] = provinces.loc['Khon Kaen']
    provinces.loc['อทัุยธาน'] = provinces.loc['Uthai Thani']
    provinces.loc['หนองบัวล าภู'] = provinces.loc['Nong Bua Lam Phu']
    provinces.loc['สตลู'] = provinces.loc['Satun']
    provinces.loc['ปทุมธาน'] = provinces.loc['Pathum Thani']
    provinces.loc['ชลบุร'] = provinces.loc['Chon Buri']
    provinces.loc['จันทบุร'] = provinces.loc['Chanthaburi']
    provinces.loc['นนทบุร'] = provinces.loc['Nonthaburi']
    provinces.loc['เพชรบุร'] = provinces.loc['Phetchaburi']
    provinces.loc['ราชบุร'] = provinces.loc['Ratchaburi']
    provinces.loc['ลพบุร'] = provinces.loc['Lopburi']
    provinces.loc['สระบุร'] = provinces.loc['Saraburi']
    provinces.loc['สิงห์บุร'] = provinces.loc['Sing Buri']
    provinces.loc['ปราจีนบุร'] = provinces.loc['Prachin Buri']
    provinces.loc['สุราษฎร์ธาน'] = provinces.loc['Surat Thani']
    provinces.loc['ชัยภูม'] = provinces.loc['Chaiyaphum']
    provinces.loc['สุรินทร'] = provinces.loc['Surin']
    provinces.loc['อุบลราชธาน'] = provinces.loc['Ubon Ratchathani']
    provinces.loc['ประจวบคีรีขันธ'] = provinces.loc['Prachuap Khiri Khan']
    provinces.loc['อุตรดิตถ'] = provinces.loc['Uttaradit']
    provinces.loc['อ านาจเจริญ'] = provinces.loc['Amnat Charoen']
    provinces.loc['อุดรธาน'] = provinces.loc['Udon Thani']
    provinces.loc['เพชรบูรณ'] = provinces.loc['Phetchabun']
    provinces.loc['บุรีรัมย'] = provinces.loc['Buri Ram']
    provinces.loc['กาฬสินธุ'] = provinces.loc['Kalasin']
    provinces.loc['สุพรรณบุร'] = provinces.loc['Suphanburi']
    provinces.loc['กาญจนบุร'] = provinces.loc['Kanchanaburi']
    provinces.loc['ล าพูน'] = provinces.loc['Lamphun']
    provinces.loc["บึงกาฬ"] = provinces.loc["Bueng Kan"]
    provinces.loc['สมทุรสำคร'] = provinces.loc['Samut Sakhon']
    provinces.loc['กรงุเทพมหำนคร'] = provinces.loc['Bangkok']
    provinces.loc['สมทุรปรำกำร'] = provinces.loc['Samut Prakan']
    provinces.loc['อำ่งทอง'] = provinces.loc['Ang Thong']
    provinces.loc['ปทมุธำน'] = provinces.loc['Pathum Thani']
    provinces.loc['สมทุรสงครำม'] = provinces.loc['Samut Songkhram']
    provinces.loc['พระนครศรอียธุยำ'] = provinces.loc['Ayutthaya']
    provinces.loc['ตำก'] = provinces.loc['Tak']
    provinces.loc['ตรำด'] = provinces.loc['Trat']
    provinces.loc['รำชบรุ'] = provinces.loc['Ratchaburi']
    provinces.loc['ฉะเชงิเทรำ'] = provinces.loc['Chachoengsao']
    provinces.loc['Mahasarakham'] = provinces.loc['Maha Sarakham']
    provinces.loc['มหำสำรคำม'] = provinces.loc['Mahasarakham']
    provinces.loc['สรุำษฎรธ์ำน'] = provinces.loc['Surat Thani']
    provinces.loc['นครรำชสมีำ'] = provinces.loc['Nakhon Ratchasima']
    provinces.loc['ปรำจนีบรุ'] = provinces.loc['Prachinburi']
    provinces.loc['ชยันำท'] = provinces.loc['Chai Nat']
    provinces.loc['กำญจนบรุ'] = provinces.loc['Kanchanaburi']
    provinces.loc['อบุลรำชธำน'] = provinces.loc['Ubon Ratchathani']
    provinces.loc['นครศรธีรรมรำช'] = provinces.loc['Nakhon Si Thammarat']
    provinces.loc['นครนำยก'] = provinces.loc['Nakhon Nayok']
    provinces.loc['ล ำปำง'] = provinces.loc['Lampang']
    provinces.loc['นรำธวิำส'] = provinces.loc['Narathiwat']
    provinces.loc['สงขลำ'] = provinces.loc['Songkhla']
    provinces.loc['ล ำพนู'] = provinces.loc['Lamphun']
    provinces.loc['อ ำนำจเจรญ'] = provinces.loc['Amnat Charoen']
    provinces.loc['หนองคำย'] = provinces.loc['Nong Khai']
    provinces.loc['หนองบวัล ำภู'] = provinces.loc['Nong Bua Lam Phu']
    provinces.loc['อดุรธำน'] = provinces.loc['Udon Thani']
    provinces.loc['นำ่น'] = provinces.loc['Nan']
    provinces.loc['เชยีงรำย'] = provinces.loc['Chiang Rai']
    provinces.loc['ก ำแพงเพชร'] = provinces.loc['Kamphaeng Phet']
    provinces.loc['พทัลงุ*'] = provinces.loc['Phatthalung']
    provinces.loc['พระนครศรอียุธยำ'] = provinces.loc['Ayutthaya']
    provinces.loc['เชีียงราย'] = provinces.loc['Chiang Rai']
    provinces.loc['เวียงจันทร์'] = provinces.loc['Unknown']  # TODO: it's really vientiane
    provinces.loc['อุทัยธาน'] = provinces.loc['Uthai Thani']
    provinces.loc['ไม่ระบุ'] = provinces.loc['Unknown']
    provinces.loc['สะแก้ว'] = provinces.loc['Sa Kaeo']
    provinces.loc['ปรมิณฑล'] = provinces.loc['Bangkok']
    provinces.loc['เพชรบรุี'] = provinces.loc['Phetchaburi']
    provinces.loc['ปตัตำนี'] = provinces.loc['Pattani']
    provinces.loc['นครสวรรค์'] = provinces.loc['Nakhon Sawan']
    provinces.loc['เพชรบุรี'] = provinces.loc['Phetchaburi']
    provinces.loc['สมุทรปราการ'] = provinces.loc['Samut Prakan']
    provinces.loc['กทม'] = provinces.loc['Bangkok']
    provinces.loc['สระบุรี'] = provinces.loc['Saraburi']
    provinces.loc['ชัยภูมิ'] = provinces.loc['Chaiyaphum']
    provinces.loc['กัมพูชา'] = provinces.loc['Unknown']  # Cambodia
    provinces.loc['มาเลเซีย'] = provinces.loc['Unknown']  # Malaysia
    provinces.loc['เรอืนจา/ทีต่อ้งขงั'] = provinces.loc['Prison']  # Prison. Currently cluster just there. might have to change later
    provinces.loc['เรอืนจาฯ'] = provinces.loc["Prison"] # Rohinja?
    #provinces.loc['เรอื นจาฯ'] = provinces.loc['Prison']

    # use the case data as it has a mapping between thai and english names
    _, cases = next(web_files("https://covid19.th-stat.com/api/open/cases", dir="json", check=False))
    cases = pd.DataFrame(json.loads(cases)["Data"])
    cases = cases.rename(columns=dict(Province="ProvinceTh", ProvinceAlt="Provinces"))
    lup_province = cases.groupby(['ProvinceId', 'ProvinceTh', 'ProvinceEn']).size().reset_index().rename({0: 'count'}, axis=1).sort_values('count', ascending=False).set_index("ProvinceEn")
    # get the proper names from provinces
    lup_province = lup_province.reset_index().rename(columns=dict(ProvinceEn="ProvinceAlt"))
    lup_province = lup_province.set_index("ProvinceAlt").join(provinces)
    lup_province = lup_province.drop(index="Unknown")
    lup_province = lup_province.set_index("ProvinceTh").drop(columns="count")

    # now bring in the thainames as extra altnames
    provinces = provinces.combine_first(lup_province)

    # bring in some appreviations
    lupurl = "https://raw.githubusercontent.com/kristw/gridmap-layout-thailand/master/src/input/provinces.csv"
    file, _ = next(web_files(lupurl, dir="json", check=False))
    abr = pd.read_csv(file)
    on_enname = abr.merge(provinces, right_index=True, left_on="enName")
    provinces = provinces.combine_first(on_enname.rename(columns=dict(thName="ProvinceAlt")).set_index("ProvinceAlt").drop(columns=["enAbbr", "enName", "thAbbr"]))
    provinces = provinces.combine_first(on_enname.rename(columns=dict(thAbbr="ProvinceAlt")).set_index("ProvinceAlt").drop(columns=["enAbbr", "enName", "thName"]))

    on_thai = abr.merge(provinces, right_index=True, left_on="thName")
    provinces = provinces.combine_first(on_thai.rename(columns=dict(enName="ProvinceAlt")).set_index("ProvinceAlt").drop(columns=["enAbbr", "thName", "thAbbr"]))
    provinces = provinces.combine_first(on_thai.rename(columns=dict(thAbbr="ProvinceAlt")).set_index("ProvinceAlt").drop(columns=["enAbbr", "enName", "thName"]))
    provinces = provinces.combine_first(on_thai.rename(columns=dict(enAbbr="ProvinceAlt")).set_index("ProvinceAlt").drop(columns=["thAbbr", "enName", "thName"]))

    # https://raw.githubusercontent.com/codesanook/thailand-administrative-division-province-district-subdistrict-sql/master/source-data.csv

    # Add in population data
    #popurl = "http://mis.m-society.go.th/tab030104.php?y=2562&p=00&d=0000&xls=y"
    popurl = "https://en.wikipedia.org/wiki/Provinces_of_Thailand"
    file, _ = next(web_files(popurl, dir="html", check=False))
    pop = pd.read_html(file)[2]
    pop = pop.join(provinces, on="Name(in Thai)").set_index("ProvinceEn").rename(columns={"Population (2019)[1]": "Population"})

    provinces = provinces.join(pop["Population"], on="ProvinceEn")

    return provinces


def get_ifr():
    file, _ = next(web_files("http://statbbi.nso.go.th/staticreport/Page/sector/EN/report/sector_01_11101_EN_.xlsx", dir="json"))
    pop = pd.read_excel(file, header=3, index_col=1)
    pop['At 0'] = pop[[f"{i} year" for i in range(1,10)]+["under 1"]].sum(axis=1)
    pop["At 10"] = pop[[f"{i} year" for i in range(10,25)]].sum(axis=1)
    pop["At 25"] =  pop[[f"{i} year" for i in range(25,46)]+["47 year"]+[f"{i} year" for i in range(47,54)]].sum(axis=1)
    pop["At 55"] =  pop[[f"{i} year" for i in range(55,65)]].sum(axis=1)
    pop["At 65"] =  pop[[f"{i} year" for i in range(65,73)]+["74 year", "74 year"]].sum(axis=1)
    pop["At 75"] = pop[[f"{i} year" for i in range(75,85)]].sum(axis=1)
    pop["At 85"] =pop[[f"{i} year" for i in range(85,101)]+["101 and over"]].sum(axis=1)   
    # from http://epimonitor.net/Covid-IFR-Analysis.htm. Not sure why pd.read_html doesn't work in this case.
    ifr = pd.DataFrame([[.002,.002,.01,.04,1.4,4.6,15]],
       columns=["At 0","At 10","At 25","At 55","At 65","At 75", "At 85"],
    ).transpose().rename(columns={0:"risk"})
    pop = pop[ifr.index]
    pop = pop.reset_index().dropna().set_index("Province").transpose()
    unpop = pop.reset_index().melt(id_vars=['index'], var_name='Province', value_name='Population').rename(columns=dict(index="Age"))
    total_pop = unpop.groupby("Province").sum().rename(columns=dict(Population="total_pop"))
    unpop = unpop.join(total_pop, on="Province").join(ifr["risk"], on="Age")
    unpop['ifr'] = unpop['Population'] / unpop['total_pop'] * unpop['risk']
    provifr = unpop.groupby("Province").sum()
    provifr = provifr.drop([p for p in provifr.index if "Region" in p] + ['Whole Kingdom'])

    # now normalise the province names
    provifr = join_provinces(provifr, on="Province")
    return provifr

PROVINCES = get_provinces()

prov_guesses = pd.DataFrame([],columns=["Province","ProvinceEn", "count"])

def get_province(prov, ignore_error=False):
    global prov_guesses
    prov = remove_prefix(prov.strip().strip(".").replace(" ", ""), "จ.")
    try:
        return PROVINCES.loc[prov]['ProvinceEn']
    except KeyError:
        try:
            close = difflib.get_close_matches(prov, PROVINCES.index)[0]
        except IndexError:
            if ignore_error:
                return None
            else:
                print(f"provinces.loc['{prov}'] = provinces.loc['x']")
                raise Exception(f"provinces.loc['{prov}'] = provinces.loc['x']")
                #continue
        proven = PROVINCES.loc[close]['ProvinceEn']  # get english name here so we know we got it
        prov_guesses = prov_guesses.append([dict(Province=prov, ProvinceEn=proven, count=1)])
        return proven

def prov_trim(p):
    return remove_suffix(remove_prefix(p, "จ.").strip(' .'), " Province")

def join_provinces(df, on):
    global prov_guesses
    joined, guesses = fuzzy_join(df, PROVINCES[["Health District Number", "ProvinceEn"]], on, True, prov_trim, "ProvinceEn", return_unmatched=True)
    if not guesses.empty:
        prov_guesses = prov_guesses.append(guesses.reset_index().rename(columns={on:"Province"})[['Province','ProvinceEn','count']]) # TODO: put in what the guess was
    return joined 

def fuzzy_join(a, b, on, assert_perfect_match=False, trim=lambda x: x, replace_on_with=None, return_unmatched=False):
    "does a pandas join but matching very similar entries"
    old_index = None
    if on not in a.columns:
        old_index = a.index.names
        a = a.reset_index()
    first = a.join(b, on=on)
    test = list(b.columns)[0]
    unmatched = first[first[test].isnull() & first[on].notna()]
    if unmatched.empty:
        second = first
    else:
        a["fuzzy_match"] = unmatched[on].map(lambda x: next(iter(difflib.get_close_matches(trim(x), b.index)), None),
                                             na_action="ignore")
        second = first.combine_first(a.join(b, on="fuzzy_match"))
        del second["fuzzy_match"]
        unmatched2 = second[second[test].isnull() & second[on].notna()]
        if assert_perfect_match:
            assert unmatched2.empty, f"Still some values left unmatched {list(unmatched2[on])}"

    unmatched_counts = pd.DataFrame()
    if return_unmatched and not unmatched.empty:
        to_keep = [test, replace_on_with] if replace_on_with is not None else [test]
        unmatched_counts = unmatched[[on]].join(second[to_keep]).value_counts().reset_index().rename(columns={0: "count"})

    if replace_on_with is not None:
        second[on] = second[replace_on_with]
        del second[replace_on_with]
    if old_index is not None:
        second = second.set_index(old_index)
    if return_unmatched:
        return second, unmatched_counts
    else:
        return second


def get_cases_by_area_type():
    dfprov, twcases = get_cases_by_prov_tweets()
    briefings, cases = get_cases_by_prov_briefings()
    cases = cases.combine_first(twcases)
    dfprov = briefings.combine_first(dfprov)  # TODO: check they aggree
    # df2.index = df2.index.map(lambda x: difflib.get_close_matches(x, df1.index)[0])
    #dfprov = dfprov.join(PROVINCES['Health District Number'], on="Province")
    dfprov = join_provinces(dfprov, on="Province")
    # Now we can save raw table of provice numbers
    export(dfprov, "cases_by_province")

    # Reduce down to health areas
    dfprov_grouped = dfprov.groupby(["Date", "Health District Number"]).sum(min_count=1).reset_index()
    dfprov_grouped = dfprov_grouped.pivot(index="Date", columns=['Health District Number'])
    dfprov_grouped = dfprov_grouped.rename(columns=dict((i, f"Area {i}") for i in DISTRICT_RANGE))
    #cols = dict((f"Area {i}", f"Cases Area {i}") for i in DISTRICT_RANGE)
    #by_area = dfprov_grouped["Cases"].groupby(['Health District Number'],axis=1).sum(min_count=1).rename(columns=cols)
    #cols = dict((f"Area {i}", f"Cases Proactive Area {i}") for i in DISTRICT_RANGE)
    by_type = dfprov_grouped.groupby(level=0, axis=1).sum(min_count=1)
    # Collapse columns to "Cases Proactive Area 13" etc
    dfprov_grouped.columns = dfprov_grouped.columns.map(' '.join).str.strip()
    by_area = dfprov_grouped.combine_first(by_type)
    by_area = by_area.combine_first(cases)  # imported, proactive total etc

    # Ensure we have all areas
    for i in DISTRICT_RANGE:
        col = f"Cases Walkin Area {i}"
        if col not in by_area:
            by_area[col] = by_area.get(col, pd.Series(index=by_area.index, name=col))
        col = f"Cases Proactive Area {i}"
        if col not in by_area:
            by_area[col] = by_area.get(col, pd.Series(index=by_area.index, name=col))
    return by_area


def get_case_details_csv():
    url = "https://data.go.th/dataset/covid-19-daily"
    file, text = next(web_files(url, dir="json", check=True))
    data = re.search(r"packageApp\.value\('meta',([^;]+)\);", text.decode("utf8")).group(1)
    apis = json.loads(data)
    links = [api['url'] for api in apis if "รายงานจำนวนผู้ติดเชื้อ COVID-19 ประจำวัน" in api['name']]
    #links = [l for l in web_links(url, ext=".csv") if "pm-" in l]
    file, _ = next(web_files(*links, dir="json"))
    if file.endswith(".xlsx"):
        cases = pd.read_excel(file)
    elif file.endswith(".csv"):
        cases = pd.read_csv(file)
    else:
        raise Exception(f"Unknown filetype for covid19daily {file}")
    cases['announce_date'] = pd.to_datetime(cases['announce_date'], dayfirst=True)
    cases['Notified date'] = pd.to_datetime(cases['Notified date'], dayfirst=True,)
    cases = cases.rename(columns=dict(announce_date="Date")).set_index("Date")
    print("Covid19daily", file, cases.reset_index().iloc[-1].to_string(header=False, index=False))
    return cases


def get_case_details_api():
    # _, cases = next(web_files("https://covid19.th-stat.com/api/open/cases", dir="json"))
    url = "https://data.go.th/api/3/action/datastore_search?resource_id=329f684b-994d-476b-91a4-62b2ea00f29f&limit=1000&offset="
    records = []

    def get_page(i, check=False):
        _, cases = next(web_files(f"{url}{i}", dir="json", check=check))
        return json.loads(cases)['result']['records']

    for i in range(0, 100000, 1000):
        data = get_page(i, False)
        if len(data) < 1000:
            data = get_page(i, True)
            if len(data) < 1000:
                break
        records.extend(data)
    # they screwed up the date conversion. d and m switched sometimes
    # TODO: bit slow. is there way to do this in pandas?
    for record in records:
        record['announce_date'] = to_switching_date(record['announce_date'])
        record['Notified date'] = to_switching_date(record['Notified date'])
    cases = pd.DataFrame(records)
    return cases


def get_cases_by_area_api():
    cases = get_case_details_csv().reset_index()
    cases["province_of_onset"] = cases["province_of_onset"].str.strip(".")
    cases = join_provinces(cases, "province_of_onset")
    case_areas = pd.crosstab(cases['Date'], cases['Health District Number'])
    case_areas = case_areas.rename(columns=dict((i, f"Cases Area {i}") for i in DISTRICT_RANGE))
    return case_areas


def get_cases_by_demographics_api():
    print("========Covid19Daily Demographics==========")

    cases = get_case_details_csv().reset_index()
    #cases = cases.rename(columns=dict(announce_date="Date"))

    #age_groups = pd.cut(cases['age'], bins=np.arange(0, 100, 10))
    # cases = get_case_details_csv().reset_index()
    labels = ["Age 0-19", "Age 20-29", "Age 30-39", "Age 40-49", "Age 50-65", "Age 66-"]
    age_groups = pd.cut(cases['age'], bins=[0, 19, 29, 39, 49, 65, np.inf], labels=labels)
    case_ages = pd.crosstab(cases['Date'], age_groups)
    #case_areas = case_areas.rename(columns=dict((i,f"Cases Area {i}") for i in DISTRICT_RANGE))

    cases['risk'].value_counts()
    risks = {}
    risks['สถานบันเทิง'] = "Entertainment"
    risks['อยู่ระหว่างการสอบสวน'] = "Investigating"  # Under investication
    risks['การค้นหาผู้ป่วยเชิงรุกและค้นหาผู้ติดเชื้อในชุมชน'] = "Proactive Search"
    risks['State Quarantine'] = 'Imported'
    risks['ไปสถานที่ชุมชน เช่น ตลาดนัด สถานที่ท่องเที่ยว'] = "Community"
    risks['Cluster ผับ Thonglor'] = "Entertainment"
    risks['ผู้ที่เดินทางมาจากต่างประเทศ และเข้า ASQ/ALQ'] = 'Imported'
    risks['Cluster บางแค'] = "Community"  # bangkhee
    risks['Cluster ตลาดพรพัฒน์'] = "Community"  #market
    risks['Cluster ระยอง'] = "Entertainment"  # Rayong
    risks['อาชีพเสี่ยง เช่น ทำงานในสถานที่แออัด หรือทำงานใกล้ชิดสัมผัสชาวต่างชาติ เป็นต้น'] = "Work"  # work with forigners
    risks['ศูนย์กักกัน ผู้ต้องกัก'] = "Prison"  # detention
    risks['คนไทยเดินทางกลับจากต่างประเทศ'] = "Imported"
    risks['สนามมวย'] = "Entertainment"  # Boxing
    risks['ไปสถานที่แออัด เช่น งานแฟร์ คอนเสิร์ต'] = "Community"  # fair/market
    risks['คนต่างชาติเดินทางมาจากต่างประเทศ'] = "Imported"
    risks['บุคลากรด้านการแพทย์และสาธารณสุข'] = "Work"
    risks['ระบุไม่ได้'] = "Unknown"
    risks['อื่นๆ'] = "Unknown"
    risks['พิธีกรรมทางศาสนา'] = "Community"  # Religous
    risks['Cluster บ่อนพัทยา/ชลบุรี'] = "Entertainment"  # gambling rayong
    risks['ผู้ที่เดินทางมาจากต่างประเทศ และเข้า HQ/AHQ'] = "Imported"
    risks['Cluster บ่อนไก่อ่างทอง'] = "Entertainment"  # cockfighting
    risks['Cluster จันทบุรี'] = "Entertainment"  # Chanthaburi - gambing?
    risks['Cluster โรงงาน Big Star'] = "Work"  # Factory
    r = {
        27: 'Cluster ชลบุรี:Entertainment',  # Chonburi - gambling
        28: 'Cluster เครือคัสเซ่อร์พีคโฮลดิ้ง (CPG,CPH):Work',
        29: 'ตรวจก่อนทำหัตถการ:Unknown',  #'Check before the procedure'
        30: 'สัมผัสผู้เดินทางจากต่างประเทศ:Contact',  # 'touch foreign travelers'
        31: "Cluster Memory 90's กรุงเทพมหานคร:Entertainment",
        32: 'สัมผัสผู้ป่วยยืนยัน:Contact',
        33: 'ปอดอักเสบ (Pneumonia):Pneumonia',
        34: 'Cluster New Jazz กรุงเทพมหานคร:Entertainment',
        35: 'Cluster มหาสารคาม:Entertainment',  # Cluster Mahasarakham
        36: 'ผู้ที่เดินทางมาจากต่างประเทศ และเข้า OQ:Imported',
        37: 'Cluster สมุทรปราการ (โรงงาน บริษัทเมทัล โปรดักส์):Work',
        38: 'สัมผัสใกล้ชิดผู้ป่วยยันยันก่อนหน้า:Contact',
        39: 'Cluster ตลาดบางพลี:Work',
        40: 'Cluster บ่อนเทพารักษ์:Community',  # Bangplee Market'
        41: 'Cluster Icon siam:Community',
        42: 'Cluster The Lounge Salaya:Entertainment',
        43: 'Cluster ชลบุรี โรงเบียร์ 90:Entertainment',
        44: 'Cluster โรงงาน standard can:Work',
        45: 'Cluster ตราด:Community',  # Trat?
        46: 'Cluster สถานบันเทิงย่านทองหล่อ:Entertainment',
        47: 'ไปยังพื้นที่ที่มีการระบาด:Community',
        48: 'Cluster สมุทรสาคร:Work',  #Samut Sakhon
        49: 'สัมผัสใกล้ชิดกับผู้ป่วยยืนยันรายก่อนหน้านี้:Contact',
        51: 'อยู่ระหว่างสอบสวน:Unknown',
        20210510.1: 'Cluster คลองเตย:Community',  # Cluster Klongtoey, 77
        20210510.2: 'ไปแหล่งชุมชน/สถานที่คนหนาแน่น:Community',  # Go to a community / crowded place, 17
        20210510.3: 'สัมผัสใกล้ชิดผู้ป่วยยืนยันก่อนหน้า:Contact',
        20210510.4: 'Cluster ชลบุรี บริษัทไดกิ้น:Work',  # Cluster Chonburi Daikin Company, 3
        20210510.5: 'ร้านอาหาร:Entertainment',  #resturant
        20210510.6: 'สัมผัสผู้ติดเชื้อยืนยัน อยู่ระหว่างสอบสวน:Contact',  # touch the infected person confirm Under investigation, 5
        20210510.7: 'สัมผัสผู้ป่วยยืนยัน อยู่ระหว่างสอบสวน:Contact',  # touch the infected person confirm Under investigation, 5
        20210510.8: 'ผู้เดินทางมาจากพื้นที่เสี่ยง กรุงเทพมหานคร:Community',  # Travelers from high-risk areas Bangkok, 2
        20210510.9: 'ไปยัง/มาจาก พื้นที่ระบาดกรุงเทพมหานครมหานคร:Community',  # to / from Epidemic area, Bangkok Metropolis, 1
        20210510.11: 'ระหว่างสอบสวน:Investigating',
        20210510.12: 'Cluster ปากช่อง:Entertainment',  # cluster pakchong https://www.bangkokpost.com/thailand/general/2103827/5-covid-clusters-in-nakhon-ratchasima - birthday party
        20210512.1: 'Cluster คลองเตย:Community',  # klongtoey cluster
        20210512.2: 'อยู่ระหว่างสอบสวนโรค:Investigating',
        20210512.3: 'อื่น ๆ:Unknown',  # Other
        20210512.4: 'Cluster จันทบุรี (ชาวกินี ):Entertainment',  # African gem merchants dining after ramandan
        20210516.0: 'Cluster เรือนจำกลางคลองเปรม:Prison',  # 894
        20210516.1: 'Cluster ตลาดสี่มุมเมือง:Community',  # 344 Four Corners Market
        20210516.2: 'Cluster สมุทรปราการ GRP Hightech:Work',  #130
        20210516.3: 'Cluster ตลาดนนทบุรี:Community',  # Cluster Talat Nonthaburi, , 85
        20210516.4: 'Cluster โรงงาน QPP ประจวบฯ:Work',  #69
        20210516.5: 'Cluster เรือนจำพิเศษธนบุรี:Prison',  #41 Cluster Special Prison Thonburi,
        20210516.6: 'Cluster จันทบุรี (ชาวกินี):Entertainment',  #26 Cluster Chanthaburi (Guinea),
        #20210516.7: 'Cluster บริษัทศรีสวัสดิ์,Work',  #16
        20210516.8: 'อื่น:Unknown',  # 10
        20210516.9: 'Cluster เรือนจำพิเศษมีนบุรี:Prison',  #5
        20210516.11: 'Cluster จนท. สนามบินสุวรรณภูมิ:Work',  #4
        20210516.12: 'สัมผัสผู้ป่วยที่ติดโควิด:Contact',  #4
    }
    for v in r.values():
        key, cat = v.split(":")
        risks[key] = cat
    risks = pd.DataFrame(risks.items(), columns=["risk", "risk_group"]).set_index("risk")
    cases_risks, unmatched = fuzzy_join(cases, risks, on="risk", return_unmatched=True)
    matched = cases_risks[["risk", "risk_group"]]
    case_risks = pd.crosstab(cases_risks['Date'], cases_risks["risk_group"])
    case_risks.columns = [f"Risk: {x}" for x in case_risks.columns]

    # dump mappings to file so can be inspected
    export(matched.value_counts().to_frame("count"), "risk_groups", csv_only=True)
    export(unmatched, "risk_groups_unmatched", csv_only=True)

    return case_risks.combine_first(case_ages)


def get_cases_by_area():
    # we will add in the tweet data for the export
    case_briefings_tweets = get_cases_by_area_type()
    case_api = get_cases_by_area_api()  # can be very wrong for the last days

    case_areas = case_briefings_tweets.combine_first(case_api)

    export(case_areas, "cases_by_area")
    
    return case_areas


def parse_tweet(tw, tweet, found, *matches):
    "if tweet contains any of matches return its text joined with comments by the same person that also match (and contain [1/2] etc)"
    if not any_in(tweet.get('text', tweet.get("comment", "")), *matches):
        return ""
    text = tw.get_tweetinfo(tweet['id']).contents['text']
    if any(text in t for t in found):
        return ""
    # TODO: ensure tweets are [1/2] etc not just "[" and by same person
    if "[" not in text:
        return text
    for t in sorted(tw.get_tweetcomments(tweet['id']).contents, key=lambda t: t['id']):
        rest = parse_tweet(tw, t, found+[text], *matches)
        if rest and rest not in text:
            text += " " + rest 
    return text


def get_tweets_from(userid, datefrom, dateto, *matches):
    "return tweets from single person that match, merging in followups of the form [1/2]. Caches to speed up"

    tw = TwitterScraper()
    filename = os.path.join("tweets", f"tweets2_{userid}.pickle")
    os.makedirs("tweets", exist_ok=True)
    try:
        with open(filename, "rb") as fp:
            tweets = pickle.load(fp)
    except (IOError, OSError, pickle.PickleError, pickle.UnpicklingError) as e:
        print(f'Error detected when attempting to load the pickle file: {e}, setting an empty \'tweets\' dictionary')
        tweets = {}
    latest = max(tweets.keys()) if tweets else None
    if latest and dateto and latest >= (datetime.datetime.today() if not dateto else dateto).date():
        return tweets
    for limit in ([50, 2000, 5000] if tweets else [5000]):
        print(f"Getting {limit} tweets")       
        for tweet in sorted(tw.get_tweets(userid, count=limit).contents, key=lambda t: t['id']):
            date = tweet['created_at'].date()
            text = parse_tweet(tw, tweet, tweets.get(date, []), *matches)
            if text:
                tweets[date] = tweets.get(date, []) + [text]

        earliest = min(tweets.keys())
        latest = max(tweets.keys())
        print(f"got tweets {earliest} to {latest} {len(tweets)}")
        if earliest <= datefrom.date():  # TODO: ensure we have every tweet in sequence?
            break
        else:
            print(f"Retrying: Earliest {earliest}")
    with open(filename, "wb") as fp:
        pickle.dump(tweets, fp)
    return tweets


def get_cases_by_prov_tweets():
    print("========RB Tweets==========")
    # These are published early so quickest way to get data
    # previously also used to get per provice case stats but no longer published

    # Get tweets
    # 2021-03-01 and 2021-03-05 are missing
    new = get_tweets_from(531202184, d("2021-04-03"), None, "Official #COVID19 update", "📍")
    #old = get_tweets_from(72888855, d("2021-01-14"), d("2021-04-02"), "Official #COVID19 update", "📍")
    old = get_tweets_from(72888855, d("2021-02-21"), None, "Official #COVID19 update", "📍")
    
    unofficial = get_tweets_from(531202184, d("2021-04-03"), None, "🔴 BREAKING: Thai health ministry reporting")
    officials = {}
    provs = {}
    breaking = {}
    for date, tweets in list(new.items())+list(old.items()):
        for tweet in tweets:
            if "RT @RichardBarrow" in tweet:
                continue
            if "Official #COVID19 update" in tweet:
                officials[date] = tweet
            elif "👉" in tweet and "📍" in tweet:
                if tweet in provs.get(date, ""):
                    continue
                provs[date] = provs.get(date, "") + " " + tweet
    for date, tweets in unofficial.items():
        for tweet in tweets:
            if "🔴 BREAKING: Thai health ministry reporting" in tweet:
                breaking[date] = tweet 

    # Get imported vs walkin totals
    df = pd.DataFrame()

    def toint(s):
        return int(s.replace(',', '')) if s else None

    for date, text in sorted(officials.items(), reverse=True):
        imported, _ = get_next_number(text, "imported", before=True, default=0)
        local, _ = get_next_number(text, "local", before=True, default=0)
        cases = imported+local
        #cases_cum, _ = get_next_number(text, "Since Jan(?:uary)? 2020")
        deaths, _ = get_next_number(text, "dead +", "deaths +")
        serious, _ = get_next_number(text, "in serious condition", before=True)
        recovered, _ = get_next_number(text, "discharged", "left care", before=True)
        hospitalised, _ = get_next_number(text, "in care", before=True)
        vent, _ = get_next_number(text, "on ventilators", before=True)
        cols = [
            "Date", 
            "Cases Imported", 
            "Cases Local Transmission", 
            "Cases", 
            "Deaths",
            "Hospitalized",
            "Recovered",
            "Hospitalized Severe",
            "Hospitalized Respirator",
        ]
        row = [date, imported, local, cases, deaths,]
        row2 = row + [hospitalised, recovered]
        if date <= d("2021-05-01").date():
            assert not any_in(row, None), f"{date} Missing data in Official Tweet {row}"
        else:
            assert not any_in(row2, None), f"{date} Missing data in Official Tweet {row}"
        row_opt = row2 + [serious, vent]
        tdf = pd.DataFrame([row_opt], columns=cols).set_index("Date")
        print(date, "Official:", tdf.to_string(index=False, header=False))
        df = df.combine_first(tdf)    

    # do unoffical tweets in no official tweet
    for date, text in breaking.items():
        if date in officials:
            continue
        numbers, _ = get_next_numbers(text, "Thai health ministry reporting")
        if not numbers:
            continue
        deaths, cases, *_ = numbers
        cols = ["Date", "Deaths", "Cases"]
        row = [date, deaths, cases]
        tdf = pd.DataFrame([row], columns=cols).set_index("Date")
        print(date, "Breaking:", tdf.to_string(index=False, header=False))
        df = df.combine_first(tdf)    

    # get walkin vs proactive by area
    walkins = {}
    proactive = {}
    for date, text in provs.items():
        if "📍" not in text:
            continue
        if "ventilators" in text:  # after 2021-05-11 start using "👉" for hospitalisation
            continue
        start, *lines = text.split("👉", 2)
        if len(lines) < 2:
            raise Exception()
        for line in lines:
            prov_matches = re.findall(r"📍([\s\w,&;]+) ([0-9]+)", line)
            prov = dict((p.strip(), toint(v)) for ps, v in prov_matches for p in re.split("(?:,|&amp;)", ps))
            if d("2021-04-08").date() == date:
                if prov["Bangkok"] == 147:  #proactive
                    prov["Bangkok"] = 47
                elif prov["Phuket"] == 3:  #Walkins
                    prov["Chumphon"] = 3
                    prov['Khon Kaen'] = 3
                    prov["Ubon Thani"] = 7
                    prov["Nakhon Pathom"] = 6
                    prov["Phitsanulok"] = 4

            label = re.findall(r'^ *([0-9]+)([^📍👉👇\[]*)', line)
            if label:
                total, label = label[0]
                #label = label.split("👉").pop() # Just in case tweets get muddled 2020-04-07
                total = toint(total)
            else:
                raise Exception(f"Couldn't find case type in: {date} {line}")
            if total is None:
                raise Exception(f"Couldn't parse number of cases in: {date} {line}")
            elif total != sum(prov.values()):
                raise Exception(f"bad parse of {date} {total}!={sum(prov.values())}: {text}")
            if "proactive" in label:
                proactive.update(dict(((date, k), v) for k, v in prov.items()))
                print(date, "Proactive:", len(prov))
                #proactive[(date,"All")] = total                                  
            elif "walk-in" in label:
                walkins.update(dict(((date, k), v) for k, v in prov.items()))
                print(date, "Walkins:", len(prov))
                #walkins[(date,"All")] = total
            else:
                raise Exception()
    # Add in missing data
    date = d("2021-03-01")
    p = {"Pathum Thani": 35, "Nonthaburi": 1}  # "All":36,
    proactive.update(((date, k), v) for k, v in p.items())
    w = {"Samut Sakhon": 19, "Tak": 3, "Nakhon Pathom": 2, "Bangkok": 2, "Chonburi": 1, "Ratchaburi": 1}  # "All":28,
    walkins.update(((date, k), v) for k, v in w.items())
                
    cols = ["Date", "Province", "Cases Walkin", "Cases Proactive"]
    rows = []
    for date, province in set(walkins.keys()).union(set(proactive.keys())):
        rows.append([date, province, walkins.get((date, province)), proactive.get((date, province))])
    dfprov = pd.DataFrame(rows, columns=cols)
    index = pd.MultiIndex.from_frame(dfprov[['Date', 'Province']])
    dfprov = dfprov.set_index(index)[["Cases Walkin", "Cases Proactive"]]
    df = df.combine_first(cum2daily(df))
    return dfprov, df


def seperate(seq, condition):
    a, b = [], []
    for item in seq:
        (a if condition(item) else b).append(item)
    return a, b


def split(seq, condition, maxsplit=0):
    "Similar to str.split except works on lists of lines. e.g. split([1,2,3,4], lambda x: x==2) -> [[1],[2],[3,4]]"
    run = []
    last = False
    splits = 0
    for i in seq:
        if (maxsplit and splits >= maxsplit) or bool(condition(i)) == last:
            run.append(i)
        else:
            splits += 1
            yield run
            run = [i]
            last = not last            
    yield run

# def nwise(iterable, n=2):                                                      
#     iters = tee(iterable, n)                                                     
#     for i, it in enumerate(iters):                                               
#         next(islice(it, i, i), None)                                               
#     return zip(*iters)   


def pairwise(lst):
    "Takes a list and turns them into pairs of tuples, e.g. [1,2,3,4] -> [[1,2],[3,4]]"
    lst = list(lst)
    return list(zip(compress(lst, cycle([1, 0])), compress(lst, cycle([0, 1]))))    

title_num = re.compile(r"([0-9]+\.(?:[0-9]+))")

def briefing_case_detail_lines(soup):
    parts = soup.find_all('p')
    parts = [c for c in [c.strip() for c in [c.get_text() for c in parts]] if c]
    maintitle, parts = seperate(parts, lambda x: "วันที่" in x)
    if not maintitle or "ผู้ป่วยรายใหม่ประเทศไทย" not in maintitle[0]:
        return
    #footer, parts = seperate(parts, lambda x: "กรมควบคุมโรค กระทรวงสาธารณสุข" in x)
    table = list(split(parts, re.compile(r"^\w*[0-9]+\.").match))
    if len(table) == 2:
        # titles at the end
        table, titles = table
        table = [titles, table]
    else:
        table.pop(0)

    # if only one table we can use camelot to get the table. will be slow but less problems
    #ctable = camelot.read_pdf(file, pages="6", process_background=True)[0].df
        
    for titles, cells in pairwise(table):
        title = titles[0].strip("(ต่อ)").strip()
        header, cells = seperate(cells, re.compile("ลักษณะผู้ติดเชื้อ").search)
        # "อยู่ระหว่างสอบสวน (93 ราย)" on 2021-04-05 screws things up as its not a province
        # have to use look behind
        thai = "[\u0E00-\u0E7Fa-zA-Z'. ]+[\u0E00-\u0E7Fa-zA-Z'.](?<!อยู่ระหว่างสอบสวน)(?<!ยู่ระหว่างสอบสวน)(?<!ระหว่างสอบสวน)"
        nl = " *\n* *"
        #nl = " *"
        nu = "(?:[0-9]+)"
        is_pcell = re.compile(rf"({thai}(?:{nl}\({thai}\))?{nl}\( *{nu} *ราย *\))")
        lines = pairwise(islice(is_pcell.split("\n".join(cells)), 1, None))  # beacause can be split over <p>
        yield title, lines


def briefing_case_detail(date, pages):

    num_people = re.compile(r"([0-9]+) *ราย")

    totals = dict()  # groupname -> running total
    all_cells = {}
    rows = []
    if date <= d("2021-02-26"):  #missing 2nd page of first lot (1.1)
        pages = []
    for soup in pages:
        for title, lines in briefing_case_detail_lines(soup):
            if "ติดเชื้อจากต่างประเทศ" in title:  # imported
                continue
            elif "การคัดกรองเชิงรุก" in title:
                case_type = "Proactive"
            elif "เดินทางมาจากต่างประเทศ" in title:
                # case_type = "Quarantine"
                continue  # just care about province cases for now
            #if re.search("(จากระบบเฝ้าระวัง|ติดเชื้อในประเทศ)", title):
            else:
                case_type = "Walkin"
            all_cells.setdefault(title, []).append(lines)
            #print(title,case_type)

            for prov_num, line in lines:
                # for prov in provs: # TODO: should really be 1. make split only split 1.
                # TODO: sometimes cells/data seperated by "-" 2021-01-03

                prov, num = prov_num.strip().split("(", 1)
                prov = get_province(prov)
                num = int(num_people.search(num).group(1))
                totals[title] = totals.get(title, 0) + num

                _, rest = get_next_numbers(line, "(?:nผล|ผลพบ)")  # "result"
                asym, rest = get_next_number(rest, "(?s)^.*(?:ไม่มีอาการ|ไมมี่อาการ|ไม่มีอาการ)", default=0, remove=True)
                sym, rest = get_next_number(rest, "(?s)^.*(?<!(?:ไม่มี|ไมมี่|ไม่มี))(?:อาการ|อาการ)", default=0, remove=True)
                unknown, _ = get_next_number(
                    rest,
                    "อยู่ระหว่างสอบสวนโรค",
                    # "อยู่ระหว่างสอบสวน",
                    "อยู่ระหว่างสอบสวน",
                    "อยู่ระหว่างสอบสวน",
                    "ไม่ระบุ",
                    default=0)
                # unknown2 = get_next_number(
                #     rest,
                #     "อยู่ระหว่างสอบสวน",
                #     "อยู่ระหว่างสอบสวน",
                #     default=0)
                # if unknown2:
                #     unknown = unknown2

                # TODO: if 1, can be by itself
                if asym == 0 and sym == 0 and unknown == 0:
                    sym, asym, unknown = None, None, None
                else:
                    assert asym + sym + unknown == num
                #print(num,prov)
                rows.append((date, prov, case_type, num, asym, sym))
    # checksum on title totals
    for title, total in totals.items():
        m = num_people.search(title)
        if not m:
            continue
        if date in [d("2021-03-19")]:  #1.1 64!=56
            continue
        assert total == int(m.group(1)), f"group total={total} instead of: {title}\n{all_cells[title]}"
    df = pd.DataFrame(rows, columns=["Date", "Province", "Case Type", "Cases", "Cases Asymptomatic", "Cases Symptomatic"]).set_index(['Date', 'Province'])

    return df


def briefing_case_types(date, pages):
    rows = []
    if date < d("2021-02-01"):
        pages = []
    for i, soup in enumerate(pages):
        text = soup.get_text()
        if "รายงานสถานการณ์" not in text:
            continue
        #cases = get_next_number(text, "ติดเชื้อจาก", before=True)
        #walkins = get_next_number(text.split("รายผู้ที่เดิน")[0], "ในประเทศ", until="ราย")
        #quarantine = get_next_number(text, "ต่างประเทศ", until="ราย", default=0)
        if date == d("2021-05-17"):
            numbers, rest = get_next_numbers(text.split("อาการหนัก")[1], "ในประเทศ")
            local, cases, imported, prison, walkins, proactive, imported2, prison2, *_ = numbers
            assert local == walkins + proactive
            assert imported == imported2
            assert prison == prison2
        else:
            numbers, rest = get_next_numbers(text, "รวม", until="รายผู้ที่เดิน")
            cases, walkins, proactive, *quarantine = numbers
            quarantine = quarantine[0] if quarantine else 0
            ports, rest = get_next_number(text, "ช่องเส้นทางธรรมชาติ", "รายผู้ที่เดินทางมาจากต่างประเทศ", before=True, default=0)
            imported = ports + quarantine
            prison, _ = get_next_number(text.split("รวม")[1], "ที่ต้องขัง", default=0, until="ราย")
        proactive += prison  # not sure if they are going to add this category going forward?

        assert cases == walkins + proactive + imported, f"{date}: briefing case types don't match"

        # hospitalisations
        numbers, rest = get_next_numbers(text, "อาการหนัก")
        if numbers:                
            severe, respirator, *_ = numbers
            hospital, _ = get_next_number(text, "ใน รพ.")
            field, _ = get_next_number(text, "รพ.สนาม")
            num, _ = get_next_numbers(text, "ใน รพ.", before=True)
            hospitalised = num[0]
            assert hospital + field == hospitalised
        else:
            hospital, field, severe, respirator, hospitalised = [None]*5

        if date < d("2021-05-18"):
            recovered, _ = get_next_number(text, "(เพ่ิมขึ้น|เพิ่มขึ้น)", until="ราย")
        else:
            # 2021-05-18 Using single infographic with 3rd wave numbers?
            numbers, _ = get_next_numbers(text, "หายป่วยแล้ว")
            cum_recovered_3rd, recovered, *_ = numbers
            if cum_recovered_3rd < recovered:
                recovered = cum_recovered_3rd
        
        assert recovered is not None

        deaths, _ = get_next_number(text, "เสียชีวิตสะสม", before=True)

        # cases by region
        # bangkok, _ = get_next_number(text, "กรุงเทพฯ และนนทบุรี")
        # north, _ = get_next_number(text, "ภาคเหนือ") 
        # south, _ = get_next_number(text, "ภาคใต้")
        # east, _ = get_next_number(text, "ภาคตะวันออก")
        # central, _ = get_next_number(text, "ภาคกลาง")
        # all_regions = north+south+east+central
        # assert hospitalised == all_regions, f"Regional hospitalised {all_regions} != {hospitalised}"

        rows.append([
            date, 
            cases, 
            walkins, 
            proactive, 
            imported, 
            hospital, 
            field, 
            severe, 
            respirator, 
            hospitalised,
            recovered,
            deaths,
        ])
        break
    df = pd.DataFrame(rows, columns=[
        "Date", 
        "Cases", 
        "Cases Walkin", 
        "Cases Proactive", 
        "Cases Imported",
        "Hospitalized Hospital",
        "Hospitalized Field",
        "Hospitalized Severe",
        "Hospitalized Respirator",
        "Hospitalized",
        "Recovered",
        "Deaths",
    ]).set_index(['Date'])
    print(f"{date.date()} Briefing Cases:", df.to_string(header=False, index=False))
    return df


NUM_OR_DASH = re.compile(r"([0-9\,\.]+|-)-?")

def parse_numbers(lst):
    return [float(i.replace(",", "")) if i != "-" else 0 for i in lst]


def briefing_province_cases(date, pages):
    if date < d("2021-01-13"):
        pages = []
    rows = {}
    for i, soup in enumerate(pages):
        #camelot.read_pdf(file,pages=' '.join([str(i) for i in [i]]))
        if "อโควิดในประเทศรายใหม่" not in str(soup):
            continue
        parts = [l.get_text() for l in soup.find_all("p")]
        parts = [l for l in parts if l]
        #parts = list(split(parts, lambda x: "รวม" in x))[-1]
        preamble, *tables = split(parts, re.compile(r"รวม\s*\(ราย\)").search)
        if len(tables) <= 1:
            continue  #Additional top 10 report. #TODO: better detection of right report
        else:
            title, parts = tables
        while parts and "รวม" in parts[0]:
            totals, *parts = parts
        parts = [c.strip() for c in NUM_OR_DASH.split("\n".join(parts)) if c.strip()]
        while True:
            if len(parts) < 9:
                # TODO: can be number unknown cases - e.g. หมายเหตุ : รอสอบสวนโรค จานวน 337 ราย
                break
            if NUM_OR_DASH.search(parts[0]):
                linenum, prov, *parts = parts
            else:
                # for some reason the line number doesn't show up? but its there in the pdf...
                break
                #prov, *parts = parts 
            numbers, parts = parts[:9], parts[9:]
            thai = prov.strip().strip(" ี").strip(" ์").strip(" ิ")
            if thai in ['กทม. และปรมิ ณฑล', 'รวมจงัหวดัอนื่ๆ(']:
                # bangkok + subrubrs, resst of thailand
                break
            prov = get_province(thai)
            numbers = parse_numbers(numbers)
            numbers = numbers[1:-1]  # last is total. first is previous days
            assert len(numbers) == 7
            for i, cases in enumerate(reversed(numbers)):
                if i > 4:  # 2021-01-11 they use earlier cols for date ranges
                    break
                olddate = date-datetime.timedelta(days=i)
                rows[(olddate, prov)] = cases + rows.get((olddate, prov), 0)  # rare case where we need to merge
                if False and olddate == date:
                    if cases > 0:
                        print(date, linenum, thai, PROVINCES["ProvinceEn"].loc[prov], cases)
                    else:
                        print("no cases", linenum, thai, *numbers)
    df = pd.DataFrame(((d, p, c) for (d, p), c in rows.items()), columns=["Date", "Province", "Cases"]).set_index(["Date", "Province"])
    assert date >= d("2021-01-13") and not df.empty, f"Briefing on {date} failed to parse cases per province"
    return df

def parse_gender(x):
    return "Male" if "ชาย" in x else "Female"

def briefing_deaths(file, date, pages):
    # Only before the 2021-04-29
    all = pd.DataFrame()
    for i, soup in enumerate(pages):
        text = soup.get_text()
        title_re = re.compile("(ผูป่้วยโรคโควดิ-19|ผู้ป่วยโรคโควิด-19)")
        if title_re.search(text):
            # Summary of locations, reasons, medium age, etc
            #tables = camelot.read_pdf(file, pages=str(i+2), process_background=True)
            #if len(tables) < 2:
            #    continue
            #df = tables[1].df
            #pcells = df[[0, 1]].itertuples()

            # all bullets with no spaces == one cell
            pre, comorbid, _, *rest = re.split(r"(•[\d\D]*?)\n\n", text)
            _, age_text, ptext, *risks = reversed(rest)
            if "วันที่ทราบผลติดเชื้อ" in ptext:
                age_text, ptext, *risks = risks
            ptext = (pre.split("\n\n")[-1] + ptext)
            ptext = re.sub("(ละ|/จังหวัด|จังหวัด|ราย)","", ptext)
            pcells = pairwise(re.split(r"(\(?\d+\)?)", ptext))
            province_count = {}
            for provinces, num in pcells:
                # len() < 2 because some stray modifier?
                text_num, rest = get_next_number(provinces, remove=True)
                provs = [p.strip("() ") for p in rest.split() if len(p) > 1 and p.strip("() ")]
                num, _ = get_next_number(num)
                if num is None and text_num is not None:
                    num = text_num
                elif num is None:
                    raise Exception(f"No number of deaths found {date}: {text}")
                province_count.update(dict((get_province(p), num) for p in provs))
            # Congenital disease / risk factor The severity of the disease
            # congenital_disease = df[2][0]  # TODO: parse?
            # Risk factors for COVID-19 infection
            # risk_factors = df[3][0]
            numbers, *_ = get_next_numbers(text, "ค่ามัธยฐานของอายุ", "ค่ากลาง อายุ", "ค่ากลางอายุ", ints=False)
            med_age, min_age, max_age, *_ = numbers
            numbers, *_ = get_next_numbers(text, "ชาย")
            male, female, *_ = numbers

            title_num, _ = get_next_numbers(text, title_re)
            day, year, deaths_title, *_ = title_num

            assert male+female == deaths_title
            # TODO: <= 2021-04-30. there is duration med, max and 7-21 days, 1-4 days, <1

            # TODO: what if they have more than one page?
            sum = \
                pd.DataFrame([[date, male + female, med_age, min_age, max_age, male, female]],
                             columns=["Date", "Deaths", "Deaths Age Median", "Deaths Age Min", "Deaths Age Max",
                                      "Deaths Male", "Deaths Female"]).set_index("Date")
            dfprov = \
                pd.DataFrame(((date, p, c) for p, c in province_count.items()),
                             columns=["Date", "Province", "Deaths"]).set_index(["Date", "Province"])
            assert male+female == dfprov['Deaths'].sum()
            print(f"{date.date()} Deaths:", len(dfprov), "|", sum.to_string(header=False, index=False))
            return all, sum, dfprov

        elif "วิตของประเทศไทย" not in text:
            continue
        orig = None
        if date <= d("2021-04-19"):
            cells = [soup.get_text()]
        else:
            # Individual case detail for death
            orig = camelot.read_pdf(file, pages=str(i+2), process_background=True)[0].df
            if len(orig.columns) != 11:
                cells = [cell for r in orig.itertuples() for cell in r[1:] if cell]
            else:
                cells = []
        if orig is None and not cells:
            raise Exception(f"Couldn't parse deaths {date}")
        elif cells:
            # Older style, not row per death
            rows = []
            for cell in cells:
                lines = [l for l in cell.split("\n") if l.strip()]
                if "รายละเอียดผู้เสีย" in lines[0]:
                    lines = lines[1:]
                rest = '\n'.join(lines)
                death_num, rest = get_next_number(rest, "รายที่", "รายที", remove=True)
                age, rest = get_next_number(rest, "อายุ", "ผู้ป่ว", remove=True)
                num_2ndwave, rest = get_next_number(rest, "ระลอกใหม่", remove=True)
                numbers, _ = get_next_numbers(rest, "")
                if age is not None and death_num is not None:
                    pass
                elif age:
                    death_num, *_ = numbers
                elif death_num:
                    age, *_ = numbers
                else:
                    death_num, age, *_ = numbers
                assert 1 < age < 110
                assert 55 < death_num < 1500
                #assert age > 20
                gender = parse_gender(cell)
                match = re.search(r"ขณะป่วย (\S*)", cell)  # TODO: work out how to grab just province
                if match:
                    prov = match.group(1).replace("จังหวัด", "")
                    province = get_province(prov)
                else:
                    # handle province by itself on a line
                    p = [get_province(word, True) for line in lines[:3] for word in line.split()]
                    p = [pr for pr in p if pr]
                    if p:
                        province = p[0]
                    else:
                        raise Exception(f"no province found for death in: {cell}")
                rows.append([float(death_num), date, gender, age, province, None, None, None, None, None])
            df = \
                pd.DataFrame(rows, columns=['death_num', "Date", "gender", "age", "Province", "nationality",
                                            "congenital_disease", "case_history", "risk_factor_sickness",
                                            "risk_factor_death"]).set_index("death_num")
            all = all.append(df, verify_integrity=True)
            continue
        elif orig is not None:  # <= 2021-04-27
            df = orig.drop(columns=[0, 10])
            df.columns = ['death_num', "gender", "nationality", "age", "Province", "congenital_disease", "case_history", "risk_factor_sickness", "risk_factor_death"]
            df['death_num'] = pd.to_numeric(df['death_num'], errors="coerce")
            df['age'] = pd.to_numeric(df['age'], errors="coerce")
            df = df.dropna(subset=["death_num"])
            df['Date'] = date
            df['gender'] = df['gender'].map(parse_gender)  # TODO: handle mispelling
            df = df.set_index("death_num")
            df = join_provinces(df, "Province")
            all = all.append(df, verify_integrity=True)
            # parts = [l.get_text() for l in soup.find_all("p")]
            # parts = [l for l in parts if l]
            # preamble, *tables = split(parts, re.compile("ปัจจัยเสี่ยงการ").search)
            # for header,lines in pairwise(tables):
            #     _, *row_pairs = split(lines, re.compile("(\d+\s*(?:ชาย|หญิง))").search)
            #     for first, rest in pairwise(row_pairs):
            #         row = ' '.join(first) + ' '.join(rest)
            #         case_num, age, *dates = get_next_numbers("")
            #         print(row)
    if all.empty:
        print(f"{date.date()}: Deaths:  0")
        sum = \
            pd.DataFrame([[date, 0, None, None, None, 0, 0]],
                         columns=["Date", "Deaths", "Deaths Age Median", "Deaths Age Min", "Deaths Age Max",
                                  "Deaths Male", "Deaths Female"]).set_index("Date")
        dfprov = pd.DataFrame(columns=["Date", "Province", "Deaths"]).set_index(["Date", "Province"])

    else:
        #print("{date.date()} Deaths: ",all.to_string(header=False, index=False))
        # calculate daily summary stats
        med_age, min_age, max_age = all['age'].median(), all['age'].min(), all['age'].max()
        g = all['gender'].value_counts()
        male, female = g.get('Male', 0), g.get('Female', 0)
        sum = \
            pd.DataFrame([[date, male+female, med_age, min_age, max_age, male, female]],
                         columns=["Date", "Deaths", "Deaths Age Median", "Deaths Age Min", "Deaths Age Max",
                                  "Deaths Male", "Deaths Female"]).set_index("Date")
        print(f"{date.date()} Deaths: ", sum.to_string(header=False, index=False))
        dfprov = all[["Date", 'Province']].value_counts().to_frame("Deaths")
    
    # calculate per provice counts
    return all, sum, dfprov


def get_cases_by_prov_briefings():
    print("========Briefings==========")
    types = pd.DataFrame(columns=["Date", ]).set_index(['Date', ])
    date_prov = pd.DataFrame(columns=["Date", "Province"]).set_index(['Date', 'Province'])
    date_prov_types = pd.DataFrame(columns=["Date", "Province", "Case Type"]).set_index(['Date', 'Province'])
    deaths = pd.DataFrame()
    url = "http://media.thaigov.go.th/uploads/public_img/source/"
    start = d("2021-01-13")  #12th gets a bit messy but could be fixed
    end = today()
    links = (f"{url}{f.day:02}{f.month:02}{f.year-1957}.pdf" for f in daterange(start, end, 1))
    for file, text in web_files(*reversed(list(links)), dir="briefings"):
        pages = parse_file(file, html=True, paged=True)
        pages = [BeautifulSoup(page, 'html.parser') for page in pages]
        date = file2date(file)

        today_types = briefing_case_types(date, pages)
        types = types.combine_first(today_types)

        case_detail = briefing_case_detail(date, pages)
        date_prov_types = date_prov_types.combine_first(case_detail)

        prov = briefing_province_cases(date, pages)

        each_death, death_sum, death_by_prov = briefing_deaths(file, date, pages)
        
        wrong_deaths_report = date in [
            d("2021-03-19"),  # 19th was reported on 18th
            d("2021-03-18"), 
            d("2021-03-17"),  # 15th and 17th no details of death
            d("2021-03-15"), 
            d("2021-02-24"),  # 02-24 infographic is image
            d("2021-02-19"),  # 02-19 death deatils is graphic (the doctor)
            d("2021-02-15"),  # no details of deaths (2)
            d("2021-02-10"),  # no details of deaths (1)
        ] or date < d("2021-02-01")  # TODO: check out why later
        ideaths, ddeaths = today_types['Deaths'], death_sum['Deaths']
        assert wrong_deaths_report or (ddeaths == ideaths).all(), f"Death details {ddeaths} didn't match total {ideaths}"

        deaths = deaths.append(each_death, verify_integrity=True)
        date_prov = date_prov.combine_first(death_by_prov)
        types = types.combine_first(death_sum)

        date_prov = date_prov.combine_first(prov)

        # Do some checks across the data
        today_total = today_types[['Cases Proactive', "Cases Walkin"]].sum().sum()
        prov_total = prov.groupby("Date").sum()['Cases'].loc[date]
        #assert prov_total and today_total
        warning = f"briefing provs={prov_total}, cases={today_total}"
        if today_total and prov_total:
            assert prov_total/today_total > 0.77, warning  # 2021-04-17 is very low but looks correct
        if today_total != prov_total:
            print(f"{date.date()} WARNING:", warning)
        # if today_total / prov_total < 0.9 or today_total / prov_total > 1.1:
        #     raise Exception(f"briefing provs={prov_total}, cases={today_total}")

        # Phetchabun                  1.0 extra
    # ขอนแกน่ 12 missing
    # ชุมพร 1 missing

    export(deaths, "deaths")

    if not date_prov_types.empty:
        symptoms = date_prov_types[["Cases Symptomatic", "Cases Asymptomatic"]]  # todo could keep province breakdown
        symptoms = symptoms.groupby(['Date']).sum()
        types = types.combine_first(symptoms)
        date_prov_types = date_prov_types[["Case Type", "Cases"]]
        date_prov_types = date_prov_types.groupby(['Date', 'Province', 'Case Type']).sum()  # we often have multiple walkin events
        date_prov_types = date_prov_types.reset_index().pivot(index=["Date", "Province"], columns=['Case Type'])
        date_prov_types.columns = [f"Cases {c}" for c in date_prov_types.columns.get_level_values(1)]
        date_prov = date_prov.combine_first(date_prov_types)

    return date_prov, types


def add_data(data, df):
    "Appends while dropping any duplicate rows"
    try:
        data = data.append(df, verify_integrity=True)
    except ValueError:
        print('detected duplicates; dropping only the duplicate rows')
        idx_names = data.index.names
        data = data.reset_index().append(df.reset_index()).drop_duplicates()
        data = data.set_index(idx_names)
    return data


def get_hospital_resources():
    print("========ArcGIS==========")

    # PUI + confirmed, recovered etc stats
#    fields = ['OBJECTID', 'ID', 'agency_code', 'label', 'agency_status', 'status', 'address', 'province', 'amphoe', 'tambol', 'latitude', 'longitude', 'level_performance', 'ministryname', 'depart', 'ShareRoom_Total', 'ShareRoom_Available', 'ShareRoom_Used', 'Private_AIIR_Total', 'Private_AIIR_Available', 'Private_AIIR_Used', 'Private_Modified_AIIR_Total', 'Private_Modified_AIIR_Available', 'Private_Modified_AIIR_Used', 'Private_Isolation_room_Total', 'Private_Isolation_room_Availabl', 'Private_Isolation_room_Used', 'Private_Cohort_ward_Total', 'Private_Cohort_ward_Available', 'Private_Cohort_ward_Used', 'Private_High_Flow_Total', 'Private_High_Flow_Available', 'Private_High_Flow_Used', 'Private_OR_negative_pressure_To', 'Private_OR_negative_pressure_Av', 'Private_OR_negative_pressure_Us', 'Private_ICU_Total', 'Private_ICU_Available', 'Private_ICU_Used', 'Private_ARI_clinic_Total', 'Private_ARI_clinic_Available', 'Private_ARI_clinic_Used', 'Volume_control_Total', 'Volume_control_Available', 'Volume_control_Used', 'Pressure_control_Total', 'Pressure_control_Available', 'Pressure_control_Used', 'Volumecontrol_Child_Total', 'Volumecontrol_Child_Available', 'Volumecontrol_Child_Used', 'Ambulance_Total', 'Ambulance_Availble', 'Ambulance_Used', 'Pills_Favipiravir_Total', 'Pills_Favipiravir_Available', 'Pills_Favipiravir_Used', 'Pills_Oseltamivir_Total', 'Pills_Oseltamivir_Available', 'Pills_Oseltamivir_Used', 'Pills_ChloroquinePhosphate_Tota', 'Pills_ChloroquinePhosphate_Avai', 'Pills_ChloroquinePhosphate_Used', 'Pills_LopinavirRitonavir_Total', 'Pills_LopinavirRitonavir_Availa', 'Pills_LopinavirRitonavir_Used', 'Pills_Darunavir_Total', 'Pills_Darunavir_Available', 'Pills_Darunavir_Used', 'Lab_PCRTest_Total', 'Lab_PCRTest_Available', 'Lab_PCRTest_Used', 'Lab_RapidTest_Total', 'Lab_RapidTest_Available', 'Lab_RapidTest_Used', 'Face_shield_Total', 'Face_shield_Available', 'Face_shield_Used', 'Cover_all_Total', 'Cover_all_Available', 'Cover_all_Used', 'ถุงมือไนไตรล์ชนิดใช้', 'ถุงมือไนไตรล์ชนิดใช้_1', 'ถุงมือไนไตรล์ชนิดใช้_2', 'ถุงมือไนไตรล์ชนิดใช้_3', 'ถุงมือไนไตรล์ชนิดใช้_4', 'ถุงมือไนไตรล์ชนิดใช้_5', 'ถุงมือยางชนิดใช้แล้ว', 'ถุงมือยางชนิดใช้แล้ว_1', 'ถุงมือยางชนิดใช้แล้ว_2', 'ถุงสวมขา_Leg_cover_Total', 'ถุงสวมขา_Leg_cover_Available', 'ถุงสวมขา_Leg_cover_Used', 'พลาสติกหุ้มคอ_HOOD_Total', 'พลาสติกหุ้มคอ_HOOD_Available', 'พลาสติกหุ้มคอ_HOOD_Used', 'พลาสติกหุ้มรองเท้า_Total', 'พลาสติกหุ้มรองเท้า_Availab', 'พลาสติกหุ้มรองเท้า_Used', 'แว่นครอบตาแบบใส_Goggles_Total', 'แว่นครอบตาแบบใส_Goggles_Availab', 'แว่นครอบตาแบบใส_Goggles_Used', 'เสื้อกาวน์ชนิดกันน้ำ_T', 'เสื้อกาวน์ชนิดกันน้ำ_A', 'เสื้อกาวน์ชนิดกันน้ำ_U', 'หมวกคลุมผมชนิดใช้แล้', 'หมวกคลุมผมชนิดใช้แล้_1', 'หมวกคลุมผมชนิดใช้แล้_2', 'เอี๊ยมพลาสติกใส_Apron_Total', 'เอี๊ยมพลาสติกใส_Apron_Available', 'เอี๊ยมพลาสติกใส_Apron_Used', 'UTM_Total', 'UTM_Available', 'UTM_Used', 'VTM_Total', 'VTM_Available', 'VTM_Used', 'Throat_Swab_Total', 'Throat_Swab_Available', 'Throat_Swab_Used', 'NS_Swab_Total', 'NS_Swab_Available', 'NS_Swab_Used', 'Surgicalmask_Total', 'Surgicalmask_Available', 'Surgicalmask_Used', 'N95_Total', 'N95_Available', 'N95_Used', 'Dr_ChestMedicine_Total', 'Dr_ChestMedicine_Available', 'Dr_ChestMedicine_Used', 'Dr_ID_Medicine_Total', 'Dr_ID_Medicine_Availble', 'Dr_ID_Medicine_Used', 'Dr_Medical_Total', 'Dr_Medical_Available', 'Dr_Medical_Used', 'Nurse_ICN_Total', 'Nurse_ICN_Available', 'Nurse_ICN_Used', 'Nurse_RN_Total', 'Nurse_RN_Available', 'Nurse_RN_Used', 'Pharmacist_Total', 'Pharmacist_Available', 'Pharmacist_Used', 'MedTechnologist_Total', 'MedTechnologist_Available', 'MedTechnologist_Used', 'Screen_POE', 'Screen_Walk_in', 'PUI', 'Confirm_mild', 'Confirm_moderate', 'Confirm_severe', 'Confirm_Recovered', 'Confirm_Death', 'GlobalID', 'region_health', 'CoverAll_capacity', 'ICU_Covid_capacity', 'N95_capacity', 'AIIR_room_capacity', 'CoverAll_status', 'Asymptomatic', 'ICUforCovidTotal', 'ICUforCovidAvailable', 'ICUforCovidUsed']
#    pui =  "https://services8.arcgis.com/241MQ9HtPclWYOzM/arcgis/rest/services/Corona_Date/FeatureServer/0/query?f=json&where=1%3D1&returnGeometry=false&spatialRel=esriSpatialRelIntersects&outFields=*&orderByFields=Date%20asc&resultOffset=0&resultRecordCount=32000&resultType=standard&cacheHint=true"

#    icu = "https://services8.arcgis.com/241MQ9HtPclWYOzM/arcgis/rest/services/Hospital_Data_Dashboard/FeatureServer/0/query?f=json&where=1%3D1&returnGeometry=false&spatialRel=esriSpatialRelIntersects&outFields=*&outStatistics=%5B%7B%22statisticType%22%3A%22sum%22%2C%22onStatisticField%22%3A%22Private_ICU_Total%22%2C%22outStatisticFieldName%22%3A%22value%22%7D%5D&resultType=standard&cacheHint=true"

    rows = []
    for page in range(0, 2000, 1000):
        every_district = f"https://services8.arcgis.com/241MQ9HtPclWYOzM/arcgis/rest/services/Hospital_Data_Dashboard/FeatureServer/0/query?f=json&where=1%3D1&returnGeometry=false&spatialRel=esriSpatialRelIntersects&outFields=*&resultOffset={page}&resultRecordCount=1000&cacheHint=true"
        file, content = next(web_files(every_district, dir="json", check=True))
        jcontent = json.loads(content)
        rows.extend([x['attributes'] for x in jcontent['features']])

    data = pd.DataFrame(rows).groupby("province").sum()
    data['Date'] = today().date()
    data['Date'] = pd.to_datetime(data['Date'])
    data = data.reset_index().set_index(["Date", "province"])
    #print("Active Cases:",data.sum().to_string(index=False, header=False))
    old = import_csv("hospital_resources")
    if old is not None:
        old = old.set_index(["Date", "province"])
        # TODO: seems to be dropping old data. Need to test
        data = add_data(old, data)
    export(data, "hospital_resources", csv_only=True)
    return data


def any_in(target, *matches):
    return any(m in target for m in matches)


def area_crosstab(df, col, suffix):
    given_2 = df.reset_index()[['Date', col+suffix, 'Health District Number']]
    given_by_area_2 = pd.crosstab(given_2['Date'], given_2['Health District Number'], values=given_2[col+suffix],  aggfunc='sum')
    given_by_area_2.columns = [f"{col} Area {c}{suffix}" for c in given_by_area_2.columns]
    return given_by_area_2


def get_vaccinations():
    folders = web_links("https://ddc.moph.go.th/dcd/pagecontent.php?page=643&dept=dcd", ext=None, match=re.compile("2564"))
    links = (l for f in folders for l in web_links(f, ext=".pdf"))
    #url = "https://ddc.moph.go.th/dcd/pagecontent.php?page=647&dept=dcd"
    # Just need the latest
    pages = ((page, file2date(f), f) for f, _ in web_files(*links, dir="vaccinations") for page in parse_file(f) if file2date(f))
    vaccinations = {}
    allocations = {}
    vacnew = {}
    shots = re.compile(r"(เข็ม(?:ที|ที่|ท่ี)\s.?(?:1|2)\s*)")
    oldhead = re.compile("(เข็มที่ 1 วัคซีน|เข็มท่ี 1 และ|เข็มที ่1 และ)")
    for page, date, file in pages:  # TODO: vaccinations are the day before I think
        if not date or date <= d("2021-01-01"):  #TODO: make go back later
            continue
        date = date-datetime.timedelta(days=1)  # TODO: get actual date from titles. maybe not always be 1 day delay
        lines = [l.strip() for l in page.split('\n') if l.strip()]
        #if date > d("2021-03-22"):
        preamble, *rest = split(lines, lambda x: shots.search(x) or oldhead.search(x))
        # if preamble and "19 รำยจังหวัดสะสม ตั้งแต่วันที่" in preamble[0]: # 2021-04-26
        #     continue
        # else:
        #     preamble, *rest = split(lines, oldhead.search)
        for headings, lines in pairwise(rest):
            shot_count = max(len(shots.findall(h)) for h in headings)
            oh_count = max(len(oldhead.findall(h)) for h in headings)
            table = {12: "new_given", 10: "given", 6: "alloc"}.get(shot_count, "old_given" if oh_count else None)
            if not table:
                continue
            added = 0
            for line in lines:
                # fix some number broken in the middle
                line = re.sub(r"(\d+ ,\d+)", lambda x: x.group(0).replace(" ", ""), line)
                area, *rest = line.split(' ', 1)
                if area == "รวม" or not rest:
                    break
                if area in ["เข็มที่", "และ"]:  # Extra heading
                    continue
                cols = [c.strip() for c in NUM_OR_DASH.split(rest[0]) if c.strip()]
                if len(cols) < 5:
                    break
                if NUM_OR_DASH.match(area):
                    thaiprov, *cols = cols
                else:
                    thaiprov = area
                prov = get_province(thaiprov)
                numbers = parse_numbers(cols)
                added += 1
                if table == "alloc":
                    allocations[(date, prov)] = numbers[3:7]
                elif table == "given":
                    if len(numbers) == 16:
                        alloc_sino, alloc_az, *numbers = numbers
                    assert len(numbers) == 14
                    assert vaccinations.get((date, prov)) in [None, numbers], f"Vac {date} {prov}|{thaiprov} repeated: {numbers} != {vaccinations.get((date,prov))}"
                    vaccinations[(date, prov)] = numbers
                elif table == "new_given":
                    assert len(numbers) == 12  # some extra "-" throwing it out. have to use camelot
                    assert vacnew.get((date, prov)) in [None, numbers], f"Vac {date} {prov}|{thaiprov} repeated: {numbers} != {vaccinations.get((date,prov))}"
                    vacnew[(date, prov)] = numbers
                elif table == "old_given":
                    alloc, target_num, given, perc, *rest = numbers
                    medical, frontline, disease, elders, riskarea, *rest = rest
                    # unknown = sum(rest) # TODO: #อยู่ระหว่ำง ระบุ กลุ่มเป้ำหมำย - In the process of specifying the target group
                    vaccinations[(date, prov)] = [given, perc, 0, 0] + \
                        [medical, 0, frontline, 0, disease, 0, elders, 0, riskarea, 0]
                    allocations[(date, prov)] = [alloc, 0, 0, 0]
            assert added > 7
            print(f"{date.date()}: {table} Vaccinations: {added}")
            continue
    df = pd.DataFrame((list(key)+value for key, value in vaccinations.items()), columns=[
        "Date",
        "Province",
        "Vac Given 1 Cum",
        "Vac Given 1 %",
        "Vac Given 2 Cum",
        "Vac Given 2 %",
        "Vac Group Medical Staff 1 Cum",
        "Vac Group Medical Staff 2 Cum",
        "Vac Group Other Frontline Staff 1 Cum",
        "Vac Group Other Frontline Staff 2 Cum",
        "Vac Group Over 60 1 Cum",
        "Vac Group Over 60 2 Cum",
        "Vac Group Risk: Disease 1 Cum",
        "Vac Group Risk: Disease 2 Cum",
        "Vac Group Risk: Location 1 Cum",
        "Vac Group Risk: Location 2 Cum",
    ]).set_index(["Date", "Province"])
    # df_new = pd.DataFrame((list(key)+value for key,value in vacnew.items()), columns=[
    #     "Date",
    #     "Province",
    #     "Vac Given 1",
    #     "Vac Given 2",
    #     "Vac Group Medical Staff 1",
    #     "Vac Group Medical Staff 2",
    #     "Vac Group Other Frontline Staff 1",
    #     "Vac Group Other Frontline Staff 2",
    #     "Vac Group Over 60 1",
    #     "Vac Group Over 60 2",
    #     "Vac Group Risk: Disease 1",
    #     "Vac Group Risk: Disease 2",
    #     "Vac Group Risk: Location 1",
    #     "Vac Group Risk: Location 2",
    # ]).set_index(["Date", "Province"])
    alloc = pd.DataFrame((list(key)+value for key, value in allocations.items()), columns=[
        "Date",
        "Province",
        "Vac Allocated Sinovac 1",
        "Vac Allocated Sinovac 2",
        "Vac Allocated AstraZeneca 1",
        "Vac Allocated AstraZeneca 2",
    ]).set_index(["Date", "Province"])
    all_vac = df.combine_first(alloc)  # TODO: pesky 2021-04-26

    # Do cross check we got the same number of allocations to vaccination
    counts = all_vac.groupby("Date").count()
    missing_data = counts[counts['Vac Allocated AstraZeneca 1'] > counts['Vac Group Risk: Location 2 Cum']]
    # 2021-04-08 2021-04-06 2021-04-05- 03-02 just not enough given yet
    missing_data = missing_data["2021-04-09": "2021-05-03"]
    # 2021-05-02 2021-05-01 - use images for just one table??
    # We will just remove this days
    all_vac = all_vac.drop(index=missing_data.index)
    # After 2021-05-08 they stopped using allocation table. But cum should now always have 77 provinces
    # TODO: only have 76 prov? something going on
    missing_data = counts[counts['Vac Group Risk: Location 2 Cum'] < 76]["2021-05-04":]
    all_vac = all_vac.drop(index=missing_data.index)
    # TODO: parse the daily vaccinations to make up for missing data in cum tables

    # Fix holes in cumulative using any dailys
    # TODO: below is wrong approach. should add daily to cum -1
    # df_daily = df.reset_index().set_index("Date").groupby("Province").apply(cum2daily)
    # df_daily.combine_first(df_new)
    # df_cum = df_daily.groupby("Province").cumsum()
    # df_cum.columns = [f"{c} Cum" for c in df_cum.columns]
    # all_vac = all_vac.combine_first(df_cum)

    export(all_vac, "vaccinations", csv_only=True)

    thaivac = all_vac.groupby("Date").sum()
    thaivac.drop(columns=["Vac Given 1 %", "Vac Given 1 %"], inplace=True)

    # Get vaccinations by district
    all_vac = join_provinces(all_vac, on="Province")
    given_by_area_1 = area_crosstab(all_vac, 'Vac Given 1', ' Cum')
    given_by_area_2 = area_crosstab(all_vac, 'Vac Given 2', ' Cum')
    thaivac = thaivac.combine_first(given_by_area_1).combine_first(given_by_area_2)

    #TODO: can get todays from - https://ddc.moph.go.th/vaccine-covid19/ or briefings

    # Need to drop any dates that are incomplete.
    # TODO: could keep allocations?
    #thaivac = thaivac.drop(index=missing_data.index)

    #thaivac = thaivac.combine_first(cum2daily(thaivac))
    #thaivac = thaivac.drop([c for c in thaivac.columns if " Cum" in c], axis=1)
    # TODO: remove cumlutive and other stats we don't want

    # TODO: only return some daily summary stats
    return thaivac


# Combine and plot
def export(df, name, csv_only=False):
    print(f"Exporting: {name}")
    df = df.reset_index()
    for c in set(list(df.select_dtypes(include=['datetime64']).columns)):
        df[c] = df[c].dt.strftime('%Y-%m-%d')
    os.makedirs("api", exist_ok=True)
    # TODO: save space by dropping nan
    # json.dumps([row.dropna().to_dict() for index,row in df.iterrows()])
    if not csv_only:
        df.to_json(
            os.path.join("api", name),
            date_format="iso",
            indent=3,
            orient="records",
        )
    df.to_csv(
        os.path.join("api", f"{name}.csv"),
        index=False 
    )


def import_csv(name):
    path = os.path.join("api", f"{name}.csv")
    if not os.path.exists(path):
        return None
    old = pd.read_csv(path)
    old['Date'] = pd.to_datetime(old['Date'])
    return old


def scrape_and_combine():
    if USE_CACHE_DATA:
        # Comment out what you don't need to run
        #situation = get_situation()
        cases_by_area = get_cases_by_area()
        #vac = get_vaccinations()
        #cases_demo = get_cases_by_demographics_api()
        #tests = get_tests_by_day()
        #tests_by_area = get_tests_by_area()
        #privpublic = get_tests_private_public()
        cases = get_cases()
        pass
    else:
        cases_by_area = get_cases_by_area()
        
        situation = get_situation()
        cases_demo = get_cases_by_demographics_api()
        
        hospital = get_hospital_resources()
        vac = get_vaccinations()

        tests = get_tests_by_day()
        tests_by_area = get_tests_by_area()
        privpublic = get_tests_private_public()
        cases = get_cases()

    print("========Combine all data sources==========")
    df = pd.DataFrame(columns=["Date"]).set_index("Date")
    for f in ['cases',  'cases_by_area', 'situation', 'tests_by_area', 'tests', 'privpublic', 'cases_demo', 'vac']:            
        if f in locals():
            df = df.combine_first(locals()[f])
    print(df)

    export(prov_guesses.groupby(["Province", "ProvinceEn"]).sum(), "prov_guesses", csv_only=True)

    if USE_CACHE_DATA:
        old = import_csv("combined")
        old = old.set_index("Date")
        df = df.combine_first(old)

        return df
    else:
        export(df, "combined", csv_only=True)
        return df


# df = df.cumsum()
AREA_LEGEND_ORDERED = [
    "1: U-N: C.Mai, C.Rai, MHS, Lampang, Lamphun, Nan, Phayao, Phrae",
    "2: L-N: Tak, Phitsanulok, Phetchabun, Sukhothai, Uttaradit",
    "3: U-C: Kamphaeng Phet, Nakhon Sawan, Phichit, Uthai Thani, Chai Nat",
    "4: M-C: Nonthaburi, P.Thani, Ayutthaya, Saraburi, Lopburi, Sing Buri, Ang Thong, N.Nayok",
    "5: L-C: S.Sakhon, Kanchanaburi, N.Pathom, Ratchaburi, Suphanburi, PKK, Phetchaburi, S.Songkhram",
    "6: E: Trat, Rayong, Chonburi, S.Prakan, Chanthaburi, Prachinburi, Sa Kaeo, Chachoengsao",
    "7: M-NE: Khon Kaen, Kalasin, Maha Sarakham, Roi Et",
    "8: U-NE: S.Nakhon, Loei, U.Thani, Nong Khai, NBL, Bueng Kan, N.Phanom, Mukdahan",
    "9: L-NE: Korat, Buriram, Surin, Chaiyaphum",
    "10: E-NE: Yasothon, Sisaket, Amnat Charoen, Ubon Ratchathani",
    "11: SE: Phuket, Krabi, Ranong, Phang Nga, S.Thani, Chumphon, N.S.Thammarat",
    "12: SW: Narathiwat, Satun, Trang, Songkhla, Pattani, Yala, Phatthalung",
    "13: C: Bangkok",
]


def human_format(num: float, pos: int) -> str:
    magnitude = 0
    while abs(num) >= 1000:
        magnitude += 1
        num /= 1000.0
    # add more suffixes if you need them
    suffix = ['', 'K', 'M', 'G', 'T', 'P'][magnitude]
    return f'{num:.1f}{suffix}'


def thaipop(num: float, pos: int) -> str:
    pp = num/69630000*100
    num = num/1000000
    return f'{num:.1f}M / {pp:.1f}%'


def thaipop2(num: float, pos: int) -> str:
    pp = num/69630000/2*100
    num = num/1000000
    return f'{num:.1f}M / {pp:.1f}%'


def rearrange(l, *first):
    "reorder a list by moving first items to the front. Can be index or value"
    l = list(l)
    result = []
    for f in first:
        if type(f) != int:
            f = l.index(f)+1
        result.append(l[f-1])
        l[f-1] = None
    return result + [i for i in l if i is not None]


FIRST_AREAS = [13, 4, 5, 6, 1]  # based on size-ish
AREA_LEGEND = rearrange(AREA_LEGEND_ORDERED, *FIRST_AREAS) + ["Prison"]
AREA_LEGEND_SIMPLE = rearrange(AREA_LEGEND_ORDERED, *FIRST_AREAS)


def topprov(df, metricfunc, valuefunc=None, name="Top 5 Provinces", num=5, other_name="Rest of Thailand"):
    "return df with columns of valuefunc for the top x provinces by metricfunc"
    # Top 5 dfcine rollouts
    # old_index = df.index.names
    valuefunc = metricfunc if valuefunc is None else valuefunc

    # Apply metric on each province by itself
    with_metric = df.reset_index().set_index("Date").groupby("Province").apply(metricfunc).rename(
        0).reset_index().set_index("Date")

    # = metricfunc(df)
    last_day = with_metric.loc[with_metric.last_valid_index()]
    top5 = last_day.nlargest(num, 0).reset_index()
    # sort data into top 5 + rest
    top5[name] = top5['Province']
    df = df.join(top5.set_index("Province")[name], on="Province").reset_index()
    if other_name:
        df[name] = df[name].fillna(other_name)
        # TODO: sum() might have to be configurable?
        df = df.groupby(["Date", name]).sum().reset_index()  # condense all the "other" fields
    # apply the value function to get all the values
    values = df.set_index(["Date", name]).groupby(name, group_keys=False).apply(valuefunc).rename(0).reset_index()
    # put the provinces into cols
    series = pd.crosstab(values['Date'], values[name], values[0], aggfunc="sum")

    cols = list(top5[name])  # in right order
    if other_name:
        return series[cols + [other_name]]
    else:
        return series[cols]


def custom_cm(cm_name: str, size: int, last_colour: str = None, flip: bool = False) -> ListedColormap:
    """Returns a ListedColorMap object built with the supplied color scheme and with the last color forced to be equal
    to the parameter passed. The flip parameter allows to reverse the colour scheme if needed.
    """
    summer = matplotlib.cm.get_cmap(cm_name)
    if flip:
        newcolors = summer(np.linspace(1, 0, size))
    else:
        newcolors = summer(np.linspace(0, 1, size))

    if last_colour:
        newcolors[size - 1, :] = matplotlib.colors.to_rgba(last_colour)  # used for unknowns (ex: 'lightgrey')

    return ListedColormap(newcolors)


def clip_dataframe(df_all: pd.DataFrame, cols: Union[str, List[str]], n_rows: int) -> pd.DataFrame:
    """Removes the last n rows in the event that they contain any NaN

    :param df_all: the pandas DataFrame containing all data
    :param cols: specify columns from which to assess presence of NaN in the last n rows
    :param n_rows: the number of rows (counting from the last row, going backwards) to evaluate whether they contain
                   any NaN and if so then delete them. This deals with (possible) data missing for the most recent data
                   updates.
    """
    # detect the number of NaN in the last n rows of the DataFrame subset (i.e. only using the columns specified)
    sum_nans = df_all[cols][-n_rows:].isna().sum(axis=1)
    index_nans = sum_nans[sum_nans > 0].index

    # remove these indices from the pandas DataFrame
    cleaned_df = df_all.drop(index=index_nans)

    return cleaned_df


def get_cycle(cmap, n=None, use_index="auto"):
    if isinstance(cmap, str):
        if use_index == "auto":
            if cmap in ['Pastel1', 'Pastel2', 'Paired', 'Accent',
                        'Dark2', 'Set1', 'Set2', 'Set3',
                        'tab10', 'tab20', 'tab20b', 'tab20c']:
                use_index = True
            else:
                use_index = False
        cmap = matplotlib.cm.get_cmap(cmap)
    if not n:
        n = cmap.N
    if use_index == "auto":
        if cmap.N > 100:
            use_index = False
        elif isinstance(cmap, LinearSegmentedColormap):
            use_index = False
        elif isinstance(cmap, ListedColormap):
            use_index = True
    if use_index:
        ind = np.arange(int(n)) % cmap.N
        return cycler("color", cmap(ind))
    else:
        colors = cmap(np.linspace(0, 1, n))
        return cycler("color", colors)


def plot_area(df: pd.DataFrame, png_prefix: str, cols_subset: Union[str, Sequence[str]], title: str,
              legends: List[str] = None, kind: str = 'line', stacked=False, percent_fig: bool = True,
              unknown_name: str = 'Unknown', unknown_total: str = None, unknown_percent=False,
              ma_days: int = None, cmap: str = 'tab20',
              reverse_cmap: bool = False, highlight: List[str] = [], 
              y_formatter: Callable[[float, int], str] = human_format, clean_end=True, 
              between: List[str] = []) -> None:
    """Creates one .png file for several time periods, showing data in absolute numbers and percentage terms.

    :param df: data frame containing all available data
    :param png_prefix: file prefix (file suffix is '.png')
    :param cols_subset: specify columns from the pandas DataFrame based on either a column name prefix or based on a
                        list of column names.
    :param title: plot title
    :param legends: legends to be used on the plots (line chart and percentage)
    :param kind: the type of plot (line chart or area chart)
    :param stacked: whether the line chart should use stacked lines
    :param percent_fig: whether the percentage chart should be included
    :param unknown_name: the column name containing data related to unknowns
    :param unknown_total: the column name (to be created) with unknown totals 
    :param unknown_percent: to include unknowns in a percentage fig if enabled
    :param ma_days: number of days used when computing the moving average
    :param cmap: the matplotlib colormap to be used
    :param reverse_cmap: whether the colormap should be reversed
    :param highlight: cols to make thicker to highlight them
    :param y_formatter: function to format y axis numbers
    :param clean_end: remove days at end if there is no data (inc unknown)
    """
    if type(cols_subset) is str:
        cols = [c for c in df.columns if str(c).startswith(cols_subset)]
    else:
        cols = cols_subset
    
    if ma_days:
        for c in cols:
            df[f'{c} (MA)'] = df[c].rolling(f'{ma_days}d').mean()
        cols = [f'{c} (MA)' for c in cols]
        ma_suffix = ' (MA)'
    else:
        ma_suffix = ''

    # try to hone in on last day of "important" data. Assume first col
    last_update = df[cols[:1]].dropna().index[-1].date().strftime('%d %b %Y')  # date format chosen: '05 May 2021'
    # last_date_excl = df[cols].last_valid_index() # last date with some data (not inc unknown)

    if unknown_total:
        if ma_days:
            df[f'{unknown_total} (MA)'] = df[unknown_total].rolling(f'{ma_days}d').mean()
        total_col = f'{unknown_total}{ma_suffix}'
        unknown_col = f'{unknown_name}{ma_suffix}'
        other_cols = set(cols)-set([unknown_col])
        df[unknown_col] = df[total_col].sub(df[other_cols].sum(axis=1), fill_value=None).clip(lower=0)  #TODO: should not be 0 when no unknown_total
        if unknown_col not in cols:
            cols = cols + [unknown_col]

    if percent_fig:
        perccols = [c for c in cols if not unknown_total or unknown_percent or c != f'{unknown_name}{ma_suffix}']
        for c in  perccols:
            df[f'{c} (%)'] = df[f'{c}'] / df[perccols].sum(axis=1) * 100
        if unknown_total and not unknown_percent:
            df[f'{unknown_name}{ma_suffix} (%)'] = 0
        perccols = [f'{c} (%)' for c in cols]
        
    title = f'{title}\n'

    if ma_days:
        title = title + f'({ma_days} day rolling average)\n'
    title += f'Last Data: {last_update}\n'
    title += 'https://djay.github.io/covidthailand'

    # if legends are not specified then use the columns names else use the data passed in the 'legends' argument
    if legends is None:
        legends = [remove_suffix(c, " (MA)") for c in cols]
    elif unknown_total and unknown_name not in legends:
        legends = legends + [unknown_name]

    if unknown_total:
        colormap = custom_cm(cmap, len(cols) + 1, 'lightgrey', flip=reverse_cmap)
    else:
        colormap = custom_cm(cmap, len(cols), flip=reverse_cmap)

    # drop any rows containing 'NA' if they are in the specified columns (=subset of all columns)
    #df_clean = clip_dataframe(df_all=df, cols=cols, n_rows=10)
    last_date_unknown = df[cols].last_valid_index()  # last date with some data (inc unknown)
    if clean_end:
        df_clean = df.loc[:last_date_unknown]
    else:
        df_clean = df

    periods = {
        'all': df_clean, 
        '1': df_clean[:'2020-06-01'],
        '2': df_clean['2020-12-12':],
        '3': df_clean['2021-04-01':],
        '30d': df_clean.last('30d')
    }

    if USE_CACHE_DATA:  #TODO: have its own switch
        periods = {key: periods[key] for key in ['2']}

    for suffix, df_plot in periods.items():
        if df_plot.empty:
            continue

        if percent_fig:
            f, (a0, a1) = plt.subplots(2, 1, gridspec_kw={'height_ratios': [3, 2]}, figsize=[20, 12])
        else:
            f, a0 = plt.subplots(figsize=[20, 12])
        #plt.rcParams["axes.prop_cycle"] = get_cycle(colormap)
        a0.set_prop_cycle(get_cycle(colormap))
        
        if y_formatter is not None:
            a0.yaxis.set_major_formatter(FuncFormatter(y_formatter))

        if kind == "area":
            df_plot.plot(ax=a0, y=cols, kind=kind, stacked=stacked)
        else:
            for c in cols:
                style = "--" if c in [f"{b}{ma_suffix}" for b in between] else None
                width = 5 if c in [f"{h}{ma_suffix}" for h in highlight] else None
                df_plot.plot(ax=a0, y=c, linewidth=width, style=style, kind=kind)
        #     a0.plot(df_plot.index, df_plot.reset_index()[c])
        # if between:
        #     a0.fill_between(x=df.index.values, y1=between[0], y2=between[1], data=df)

        a0.set_title(label=title)
        a0.legend(labels=legends)

        if unknown_total:
            a0.set_ylabel(unknown_total)
        a0.xaxis.label.set_visible(False)

        if percent_fig:
            df_plot.plot(ax=a1, y=perccols, kind='area', colormap=colormap, legend=False)
            a1.set_ylabel('Percent')
            a1.xaxis.label.set_visible(False)

        plt.tight_layout()
        plt.savefig(os.path.join("outputs", f'{png_prefix}_{suffix}.png'))
        plt.close()

    return None


def save_plots(df: pd.DataFrame) -> None:
    print('======== Generating Plots ==========')

    # matplotlib global settings
    matplotlib.use('AGG')
    plt.style.use('seaborn-whitegrid')
    plt.rcParams.update({'font.size': 16})
    plt.rc('legend', **{'fontsize': 14})

    # create directory if it does not exists
    pathlib.Path('./outputs').mkdir(parents=True, exist_ok=True)

    # Computed data
    # TODO: has a problem if we have local transmission but no proactive
    # TODO: put somewhere else
    walkins = pd.DataFrame(df["Cases Local Transmission"] - df["Cases Proactive"], columns=['Cases Walkin'])
    df = df.combine_first(walkins)

    cols = ['Tests XLS', 'Tests Public', 'Tested PUI', 'Tested PUI Walkin Public', ]
    legends = ['Tests Performed (All)', 'Tests Performed (Public)', 'PUI', 'PUI (Public)', ]
    plot_area(df=df, png_prefix='tests', cols_subset=cols,
              title='Thailand PCR Tests and PUI (totals exclude some proactive testing)', legends=legends,
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    cols = ['Tested Cum',
            'Tested PUI Cum',
            'Tested Not PUI Cum',
            'Tested Proactive Cum',
            'Tested Quarantine Cum',
            'Tested PUI Walkin Private Cum',
            'Tested PUI Walkin Public Cum']
    plot_area(df=df, png_prefix='tested_pui', cols_subset=cols,
              title='PCR Tests and PUI in Thailand (excludes some proactive test)',
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    ###############
    # Positive Rate
    ###############
    df["Positivity PUI"] = df["Cases"] / df["Tested PUI"] * 100
    df["Positivity Public"] = df["Pos Public"] / df["Tests Public"] * 100
    df["Positivity Cases/Tests"] = df["Cases"] / df["Tests XLS"] * 100
    df["Positivity Public+Private"] = (df["Pos XLS"] / df["Tests XLS"] * 100)
    df['Positivity Walkins/PUI'] = df['Cases Walkin'] / df['Tested PUI'] * 100
    df['Positive Rate Private'] = (df['Pos Private'] / df['Tests Private']) * 100
    df['Cases per PUI3'] = df['Cases'] / df['Tested PUI'] / 3.0 * 100
    df['Cases per Tests'] = df['Cases'] / df['Tests XLS'] * 100

    cols = [
            'Positivity Public+Private',
            'Cases per Tests',
            'Cases per PUI3',
            'Positive Rate Private'
    ]
    legends = [
        'Positive Rate: Share of PCR tests that are positive ',
        'Share of PCR tests that have Covid',
        'Share of PUI*3 that have Covid',
        'Share of Private PCR tests that are positive'
    ]
    plot_area(df=df, png_prefix='positivity', cols_subset=cols,
              title='Positive Rate: Is enough testing happening?', legends=legends,
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10',
              highlight=['Positivity Public+Private'])

    df['PUI per Case'] = df['Tested PUI'].divide(df['Cases'])
    df['PUI3 per Case'] = df['Tested PUI']*3 / df['Cases']
    df['PUI3 per Walkin'] = df['Tested PUI']*3 / df['Cases Walkin']
    df['PUI per Walkin'] = df['Tested PUI'].divide(df['Cases Walkin'])
    df['Tests per case'] = df['Tests XLS'] / df['Cases']
    df['Tests per positive'] = df['Tests XLS'] / df['Pos XLS']

    cols = ['Tests per positive', 'Tests per case', 'PUI per Case', 'PUI3 per Case', 'PUI per Walkin']
    legends = [
            'PCR Tests per Positive',
            'PCR Tests per Case',
            'PUI per Case',
            'PUI*3 per Case',
            'PUI per Walkin Case',
        ]
    plot_area(df=df, png_prefix='tests_per_case', cols_subset=cols,
              title='Thailand Tests per Confirmed Case', legends=legends,
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    cols = ['Positivity Cases/Tests',
            'Positivity Public',
            'Positivity PUI',
            'Positive Rate Private',
            'Positivity Public+Private']
    legends = [
            'Confirmed Cases / Tests Performed (Public)',
            'Positive Results / Tests Performed (Public)',
            'Confirmed Cases / PUI',
            'Positive Results / Tests Performed (Private)',
            'Positive Results / Tests Performed (All)',
    ]
    plot_area(df=df, png_prefix='positivity_all', cols_subset=cols,
              title='Positive Rate', legends=legends,
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    ########################
    # Public vs Private
    ########################
    df['Tests Private Ratio'] = (df['Tests Private'] / df['Tests Public']).rolling('7d').mean()
    df['Tests Positive Private Ratio'] = (df['Pos Private'] / df['Pos Public']).rolling('7d').mean()
    df['Positive Rate Private Ratio'] = (df['Pos Private'] / (df['Tests Private']) / (df['Pos Public'] / df['Tests Public'])).rolling('7d').mean()
    df['PUI Private Ratio'] = (df['Tested PUI Walkin Private'] / df['Tested PUI Walkin Public']).rolling('7d').mean()
    cols = ['Tests Private Ratio',
            'Tests Positive Private Ratio',
            'PUI Private Ratio',
            'Positive Rate Private Ratio']
    plot_area(df=df, png_prefix='tests_private_ratio', cols_subset=cols, title='Testing Private Ratio',
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    ##################
    # Test Plots
    ##################
    cols = ['Cases',
            'Pos Public',
            'Pos XLS']
    legends = ['Confirmed Cases',
               'Positive Test Results (Public)',
               'Positive Test Results (All)']
    plot_area(df=df, png_prefix='cases', cols_subset=cols,
              title='Positive Test results compared to Confirmed Cases', legends=legends,
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    cols = ['Cases',
            'Pos Area',
            'Pos XLS',
            'Pos Public',
            'Pos Private',
            'Pos']
    plot_area(df=df, png_prefix='cases_all', cols_subset=cols,
              title='Positive Test results compared to Confirmed Cases',
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    cols = ['Cases Imported', 'Cases Walkin', 'Cases Proactive']
    plot_area(df=df, png_prefix='cases_types', cols_subset=cols, 
              title='Thailand Covid Cases by Where Tested',
              legends=["Quarantine (Imported)", "Hospital (Walk-ins/Traced)", "Mobile Community Testing/Prisons (Proactive)"],
              unknown_name='Cases Unknown', unknown_total='Cases',
              kind='area', stacked=True, percent_fig=False, ma_days=7, cmap="viridis")

    cols = ['Cases Symptomatic', 'Cases Asymptomatic']
    plot_area(df=df, png_prefix='cases_sym', cols_subset=cols, title='Thailand Covid Cases by Symptoms',
              unknown_name='Cases Symptomatic Unknown', unknown_total='Cases',
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab10')

    # cols = ['Cases Imported','Cases Walkin', 'Cases Proactive', 'Cases Unknown']
    # plot_area(df=df, png_prefix='cases_types_all', cols_subset=cols, title='Thailand Covid Cases by Test Type',
    #           kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab10')

    # Thailand Covid Cases by Age
    plot_area(df=df, png_prefix='cases_ages', cols_subset='Age', title='Thailand Covid Cases by Age',
              unknown_name='Unknown', unknown_total='Cases', unknown_percent=False,
              kind='area', stacked=True, percent_fig=True, ma_days=7, cmap='summer', reverse_cmap=True)

    # Thailand Covid Cases by Risk
    cols = [c for c in df.columns if str(c).startswith("Risk: ")]
    cols = rearrange(cols, "Risk: Imported", "Risk: Pneumonia", "Risk: Community", "Risk: Contact", "Risk: Work", "Risk: Entertainment", "Risk: Proactive Search", "Risk: Unknown")
    plot_area(df=df, png_prefix='cases_causes', cols_subset=cols, title='Thailand Covid Cases by Risk',
              unknown_name='Risk: Investigating', unknown_total='Cases',
              kind='area', stacked=True, percent_fig=True, ma_days=7, cmap='tab20')

    ##########################
    # Tests by area
    ##########################
    plt.rc('legend', **{'fontsize': 12})

    cols = rearrange([f'Tests Area {area}' for area in DISTRICT_RANGE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='tests_area', cols_subset=cols[0],
              title='PCR Tests by Health District (excludes proactive & private tests)', legends=AREA_LEGEND_SIMPLE,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20')

    cols = rearrange([f'Pos Area {area}' for area in DISTRICT_RANGE_SIMPLE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='pos_area', cols_subset=cols,
              title='PCR Positive Test Results by Health District (excludes proactive & private tests)',
              legends=AREA_LEGEND_SIMPLE,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20')

    for area in DISTRICT_RANGE_SIMPLE: 
        df[f'Tests Area {area} (i)'] = df[f'Tests Area {area}'].interpolate(limit_area="inside")
    test_cols = [f'Tests Area {area} (i)' for area in DISTRICT_RANGE_SIMPLE]
    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Tests Daily {area}'] = (
            df[f'Tests Area {area} (i)']
            / df[test_cols].sum(axis=1)
            * df['Tests']
        )
    cols = rearrange([f'Tests Daily {area}' for area in DISTRICT_RANGE_SIMPLE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='tests_area_daily', cols_subset=cols,
              title='PCR Tests by Thailand Health District (excludes some proactive tests)', legends=AREA_LEGEND_SIMPLE,
              #unknown_name='Unknown District', unknown_total='Tests',
              kind='area', stacked=True, percent_fig=False, ma_days=7, cmap='tab20')

    for area in DISTRICT_RANGE_SIMPLE: 
        df[f'Pos Area {area} (i)'] = df[f'Pos Area {area}'].interpolate(limit_area="inside")
    pos_cols = [f'Pos Area {area} (i)' for area in DISTRICT_RANGE_SIMPLE]
    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Pos Daily {area}'] = (
            df[f'Pos Area {area} (i)']
            / df[pos_cols].sum(axis=1)
            * df['Pos']
        )
    cols = rearrange([f'Pos Daily {area}' for area in DISTRICT_RANGE_SIMPLE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='pos_area_daily', 
              cols_subset=cols, legends=AREA_LEGEND_SIMPLE,
              title='Positive PCR Tests by Thailand Health District (excludes some proactive tests)', 
              #unknown_name='Unknown District', unknown_total='Pos',
              kind='area', stacked=True, percent_fig=False, ma_days=7, cmap='tab20')

    # Workout positivity for each area as proportion of positivity for that period
    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Positivity {area}'] = (
            df[f'Pos Area {area} (i)'] / df[f'Tests Area {area} (i)'] * 100
        )
    cols = [f'Positivity {area}' for area in DISTRICT_RANGE_SIMPLE]
    df['Total Positivity Area'] = df[cols].sum(axis=1)
    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Positivity {area}'] = (
            df[f'Positivity {area}']
            / df['Total Positivity Area']
            * df['Positivity Public+Private']
        )
    plot_area(df=df, png_prefix='positivity_area', 
              cols_subset=rearrange(cols, *FIRST_AREAS), legends=AREA_LEGEND_SIMPLE,
              title='Positive Rate by Health Area in proportion to Thailand positive rate '
                    '(excludes some proactive tests)',
              #unknown_name='Unknown District', unknown_total='Positivity Public+Private',
              kind='area', stacked=True, percent_fig=True, ma_days=7, cmap='tab20')

    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Positivity Daily {area}'] = df[f'Pos Daily {area}'] / df[f'Tests Daily {area}'] * 100
    cols = [f'Positivity Daily {area}' for area in DISTRICT_RANGE_SIMPLE]
    topcols = df[cols].sort_values(by=df[cols].last_valid_index(), axis=1, ascending=False).columns[:5]
    legend = rearrange(AREA_LEGEND_ORDERED, *[cols.index(c)+1 for c in topcols])[:5]
    plot_area(df=df, png_prefix='positivity_area_unstacked', 
              cols_subset=topcols, legends=legend,
              title='Health Districts with the highest Positive Rate', 
              kind='line', stacked=False, percent_fig=False, ma_days=7, cmap='tab10')

    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Cases/Tests {area}'] = (
            df[f'Cases Area {area}'] / df[f'Tests Area {area}'] * 100
        )
    cols = [f'Cases/Tests {area}' for area in DISTRICT_RANGE_SIMPLE]
    plot_area(df=df, png_prefix='casestests_area_unstacked', 
              cols_subset=rearrange(cols, *FIRST_AREAS), legends=AREA_LEGEND_SIMPLE,
              title='Health Districts with the highest Cases/Tests (excludes some proactive tests)',
              kind='area', stacked=False, percent_fig=False, ma_days=None, cmap='tab20')

    #########################
    # Case by area plots
    #########################
    cols = rearrange([f'Cases Area {area}' for area in DISTRICT_RANGE]+['Cases Imported'], *FIRST_AREAS)
    plot_area(df=df, png_prefix='cases_areas', 
              cols_subset=cols, legends=AREA_LEGEND + ['Imported Cases'],
              title='Thailand Covid Cases by Health District',
              unknown_name="Unknown District", unknown_total="Cases",
              kind='area', stacked=True, percent_fig=False, ma_days=7, cmap='tab20')

    cols = rearrange([f'Cases Walkin Area {area}' for area in DISTRICT_RANGE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='cases_areas_walkins', cols_subset=cols,
              title='Thailand "Walk-in" Covid Cases by Health District', legends=AREA_LEGEND,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20')

    cols = rearrange([f'Cases Proactive Area {area}' for area in DISTRICT_RANGE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='cases_areas_proactive', cols_subset=cols,
              title='Thailand "Proactive" Covid Cases by Health District', legends=AREA_LEGEND,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20')

    for area in DISTRICT_RANGE_SIMPLE:
        df[f'Case-Pos {area}'] = (
            df[f'Cases Area {area}'] - df[f'Pos Area {area}']
        )
    cols = [f'Case-Pos {area}' for area in DISTRICT_RANGE_SIMPLE]
    plot_area(df=df, png_prefix='cases_from_positives_area', 
              cols_subset=rearrange(cols, *FIRST_AREAS), legends=AREA_LEGEND_SIMPLE,
              title='Which Health Districts have more cases than positive results?', 
              kind='area', stacked=False, percent_fig=False, ma_days=None, cmap='tab20')

    #######################
    # Hospital plots
    #######################
    cols_delayed = ["Hospitalized", "Recovered", "Hospitalized Severe", "Hospitalized Respirator", "Hospitalized Field"]

    # TODO: only do for last day? Should do unknown instead? or just not show until we have all the data?
    df[cols_delayed] = df[cols_delayed].ffill()

    df["Hospitalized Severe"] = df["Hospitalized Severe"].sub(df["Hospitalized Respirator"], fill_value=0)
    non_split = df[["Hospitalized Severe", "Hospitalized Respirator", "Hospitalized Field"]].sum(skipna=False, axis=1)

    # sometimes we deaths and cases but not the rest so fillfoward.
    df["Hospitalized Hospital"] = df["Hospitalized"].sub(non_split, fill_value=0)
    cols = ["Hospitalized Respirator", "Hospitalized Severe", "Hospitalized Hospital", "Hospitalized Field"]
    legends = ['On Respirator', 'Severe Case', 'Hospitalised Other', 'Field Hospital']
    plot_area(df=df, png_prefix='cases_active', cols_subset=cols,
              title='Thailand Active Covid Cases\n(Severe, Field, and Respirator only available from '
                    '2021-04-24 onwards)',
              legends=legends,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab10')

    # show cumulitive deaths, recoveres and hospitalisations (which should all add up to cases)
    df['Recovered since 2021-04-01'] = df['2021-04-14':]['Recovered'].cumsum()
    df['Died since 2021-04-01'] = df['2021-04-01':]['Deaths'].cumsum()
    df['Cases since 2021-04-01'] = df['2021-04-01':]['Cases'].cumsum()
    df['Other Active Cases'] = \
        df['Cases since 2021-04-01'].sub(non_split, fill_value=0).sub(df['Recovered since 2021-04-01'], fill_value=0)

    cols = [
        'Died since 2021-04-01',
        'Hospitalized Respirator',
        'Hospitalized Severe',
        'Other Active Cases',
        'Hospitalized Field',
        'Recovered since 2021-04-01',
    ]
    legends = [
            'Deaths from cases since 1st April',
            'On Ventilator',
            'In severe condition',
            'In Hospital',
            'In Field Hospital',
            'Recovered from cases since 1st April'
    ]
    plot_area(df=df, png_prefix='cases_cumulative', cols_subset=cols,
              title='Current outcome of Covid Cases since 1st April 2021', legends=legends,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab10')


    ####################
    # Vaccines
    ####################
    cols = [c for c in df.columns if str(c).startswith('Vac Group')]

    def clean_vac_leg(c):
        return c.replace(' Cum', '').replace('Vac Group', '').replace('1', 'Dose 1').replace('2', 'Dose 2')

    cols.sort(key=lambda c: clean_vac_leg(c)[-1] + clean_vac_leg(c))  # put 2nd shot at end

    legends = [clean_vac_leg(c) for c in cols]
    df_vac_groups = df['2021-02-16':][cols].interpolate(limit_area="inside")
    plot_area(df=df_vac_groups, png_prefix='vac_groups', cols_subset=cols,
              title='Thailand Vaccinations by Groups\n(% of 2 doses per Thai population)', legends=legends,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='Set3',
              y_formatter=thaipop2)

    cols = rearrange([f'Vac Given 1 Area {area} Cum' for area in DISTRICT_RANGE_SIMPLE], *FIRST_AREAS)
    df_vac_areas_s1 = df['2021-02-16':][cols].interpolate()
    plot_area(df=df_vac_areas_s1, png_prefix='vac_areas_s1', cols_subset=cols,
              title='Thailand Vaccinations (1st Shot) by Health District\n(% per population)', legends=AREA_LEGEND_SIMPLE,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20',
              y_formatter=thaipop)

    cols = rearrange([f'Vac Given 2 Area {area} Cum' for area in DISTRICT_RANGE_SIMPLE], *FIRST_AREAS)
    df_vac_areas_s2 = df['2021-02-16':][cols].interpolate()
    plot_area(df=df_vac_areas_s2, png_prefix='vac_areas_s2', cols_subset=cols,
              title='Thailand Fully Vaccinated (2nd Shot) by Health District\n(% population full vaccinated)',
              legends=AREA_LEGEND_SIMPLE,
              kind='area', stacked=True, percent_fig=False, ma_days=None, cmap='tab20',
              y_formatter=thaipop)

    # Top 5 vaccine rollouts
    vac = import_csv("vaccinations")
    vac['Date'] = pd.to_datetime(vac['Date'])
    vac = vac.set_index('Date')
    vac = vac.join(PROVINCES['Population'], on='Province')
    top5 = vac.pipe(topprov, lambda df: df['Vac Given 2 Cum'] / df['Population'] * 100)

    cols = top5.columns.to_list()
    plot_area(df=top5, png_prefix='vac_top5_full', cols_subset=cols,
              title='Top 5 Thai Provinces Closest to Fully Vaccinated',
              kind='area', stacked=False, percent_fig=False, ma_days=None, cmap='tab20',
              )

    #######################
    # Cases by provinces
    #######################
    def trendline(data: pd.DataFrame, order: int = 1) -> float:
        # simulate dates with monotonic inc numbers
        dates = range(0, len(data.index.values))
        coeffs = np.polyfit(dates, list(data), order)
        slope = coeffs[-2]
        return float(slope)

    def increasing(adf: pd.DataFrame) -> pd.DataFrame:
        return adf["Cases"].rolling(3).mean().rolling(3).apply(trendline)

    def cases_ma(adf: pd.DataFrame) -> pd.DataFrame:
        return adf["Cases"].rolling(3).mean()

    def decreasing(adf: pd.DataFrame) -> pd.DataFrame:
        return 1/increasing(adf)

    def cases_ma_7(adf: pd.DataFrame) -> pd.DataFrame:
        return adf["Cases"]

    cases = import_csv("cases_by_province").set_index(["Date", "Province"])

    top5 = cases.pipe(topprov, increasing, cases_ma, name="Province Cases (3d MA)", other_name=None, num=5)
    cols = top5.columns.to_list()
    plot_area(df=top5, png_prefix='cases_prov_increasing', cols_subset=cols,
              title='Provinces with Cases Trending Up\nin last 30 days (using 3 days rolling average)',
              kind='line', stacked=False, percent_fig=False, ma_days=None, cmap='tab10')

    top5 = cases.pipe(topprov, decreasing, cases_ma, name="Province Cases (3d MA)", other_name=None, num=5)
    cols = top5.columns.to_list()
    plot_area(df=top5, png_prefix='cases_prov_decreasing', cols_subset=cols,
              title='Provinces with Cases Trending Down\nin last 30 days (using 3 days rolling average)',
              kind='line', stacked=False, percent_fig=False, ma_days=None, cmap='tab10')

    top5 = cases.pipe(topprov, cases_ma_7, name="Province Cases", other_name="Other Provinces", num=6)
    cols = top5.columns.to_list()
    plot_area(df=top5, png_prefix='cases_prov_top', cols_subset=cols,
              title='Provinces with Most Cases',
              kind='line', stacked=False, percent_fig=False, ma_days=None, cmap='tab10')


    # TODO: work out based on districts of deaths / IFR for that district
    ifr = get_ifr()
    cases = cases.join(ifr[['ifr', 'Population', 'total_pop']], on="Province")
    cases['Deaths'] = cases['Deaths'].fillna(0)
    cases = cases.groupby("Province").apply(lambda df: df.assign(deaths_ma=df["Deaths"].rolling(7, min_periods=1).mean() ))
    
    cases["Infections Estimate"] = cases['Deaths'] / (cases['ifr']/100)
    cases["Infections Estimate (MA)"] = cases['deaths_ma'] / (cases['ifr']/100)
    cases_est = cases.groupby(["Date"]).sum()

    # TODO: work out unknown deaths and use whole thailand IFR for them
    # cases_est['Deaths Unknown'] = (df['Deaths'] - cases_est['Deaths']) / ifr['ifr']['Whole Kingdom'] * 100

    cases_est["Infections Estimate"] = cases_est["Infections Estimate"].shift(-14)
    cases_est["Infections Estimate (MA)"] = cases_est["Infections Estimate (MA)"].shift(-14)
    cases_est = cases_est.rename(columns=dict(Deaths="Deaths prov sum"))
    cases_est = cases_est.join(df['Deaths'], on="Date")
    cases_est['Cases (MA)'] = cases_est['Cases'].rolling("7d").mean()
    cases_est["Infections Estimate Simple"] = cases_est["Deaths"].shift(-14) / 0.0054
    cols = ["Cases (MA)", "Infections Estimate (MA)", "Infections Estimate", "Cases",]
    legend = ["Cases (7d moving avg.)", "Lower Estimate of Infections (7d moving avg.)", "Lower Estimate of Infections", "Cases"]
    title = \
"""Thailand Confirmed Covid Cases vs Estimate of Infections based on Deaths
Estimate of Infections = (Deaths - 14days)/(Province Infection Fatality Rate)
(DISCLAIMER: estimate is simple and probably lower than reality. see site below for more details on this model)"""
    plot_area(df=cases_est, png_prefix='cases_infections_estimate', cols_subset=cols,
              title=title, 
              legends=legend,
              kind='line', stacked=False, percent_fig=False, ma_days=None, cmap='tab10',
              between=["Infections Estimate", "Cases", ])


    ####################
    # Deaths
    ####################

    # predict median age of death based on population demographics
    


    df['Deaths Age Median (MA)'] = df['Deaths Age Median'].rolling('7d').mean()
    cols = ['Deaths Age Median (MA)', 'Deaths Age Max', 'Deaths Age Min']
    plot_area(df=df, png_prefix='deaths_age', cols_subset=cols, title='Thailand Covid Death Age Range',
              kind='line', stacked=False, percent_fig=False, ma_days=None, cmap='tab10',
              highlight=['Deaths Age Median (MA)'], between=['Deaths Age Max', 'Deaths Age Min'])

    cols = rearrange([f'Deaths Area {area}' for area in DISTRICT_RANGE], *FIRST_AREAS)
    plot_area(df=df, png_prefix='deaths_by_area', cols_subset=cols,
              title='Thailand Covid Deaths by health District', legends=AREA_LEGEND,
              kind='area', stacked=True, percent_fig=True, ma_days=7, cmap='tab20')



if __name__ == "__main__":

    USE_CACHE_DATA = os.environ.get('USE_CACHE_DATA', False) == 'True' and \
                     os.path.exists(os.path.join('api', 'combined.csv'))
    print(f'\n\nUSE_CACHE_DATA = {USE_CACHE_DATA}\nCHECK_NEWER = {CHECK_NEWER}\n\n')

    df = scrape_and_combine()
    save_plots(df)
